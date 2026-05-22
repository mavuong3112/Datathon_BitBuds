"""Build 2-slide PPTX for category 1020 by cloning layout from 1030_slides.pptx."""

from __future__ import annotations

import copy
import shutil
from pathlib import Path

from pptx import Presentation

DATA_ROOT = Path(__file__).resolve().parent
TEMPLATE = DATA_ROOT / "1030_slides.pptx"
OUT_DIR = DATA_ROOT / "outputs" / "category_1020_business" / "slides"
OUT_PPTX = OUT_DIR / "category_1020_SCIS_slides.pptx"


def _delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def _replace_in_slide(slide, replacements: list[tuple[str, str]]) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        new_text = text
        for old, new in replacements:
            if old in new_text:
                new_text = new_text.replace(old, new)
        if new_text != text:
            shape.text_frame.text = new_text


def _set_shape_text(slide, index: int, text: str) -> None:
    sh = slide.shapes[index]
    if sh.has_text_frame:
        sh.text_frame.text = text


def _apply_slide1(sc) -> None:
    """SITUATION & CHALLENGES — layout from 1030 slide 2."""
    reps = [
        (
            "SITUATION & CHALLENGES — CATEGORY 1030 (VĂN PHÒNG / MẶT BẰNG)",
            "SITUATION & CHALLENGES — CATEGORY 1020 (NHÀ Ở)",
        ),
        ("93,7%", "1,51M"),
        ("Supply Let vs\n6,3% Sell", "listings catalog\n(~2,5× 1010)"),
        ("92,2%", "20,5%"),
        ("other_interaction\n(dominant B2B)", "CVR let /\n19,3% sell"),
        ("Thiếu hụt cung\n30–50 m²", "8,6% contact-day\nlet / sell"),
        ("76,1%", "6.187"),
        ("Supply là Agent", "HQ underexposed\n(~2,7× 1010)"),
        ("4,7%", "~70%"),
        ("legal_status\nNULL", "supply TP.HCM"),
        ("1,6×", "10,0%"),
        ("Non-login vs Login\nPR (bẫy đo)", "snapshot contact\nrate platform"),
        ("Thị trường\nB2B dual", "Thị trường\nBán nhà scale"),
        (
            "Category duy nhất có Let lẫn Sell đáng kể: startup/SME (30–100 m²) vs doanh nghiệp (>200 m²).\n"
            "Positive 52,5% · Revenue CVR 4,1% — nhu cầu B2B ổn định, khác residential browsing.",
            "1.507.864 listings (~2,5× cat 1010) · 72% sell / 28% let — transaction-heavy house market.\n"
            "house_type 100% · floors 61% · width 76% · ~270k login users (sample 6–8%).",
        ),
        ("Địa lý\n& BCG", "Let vs Sell\nđảo chiều"),
        (
            "Supply tập trung TP.HCM + Hà Nội; tier-2 (Đà Nẵng, Bình Dương…) ít tin nhưng CVR cao → Stars.\n"
            "C5: tỉnh tier-2 outperform trên mỗi listing — cơ hội seller acquisition ngoài HCM+HN.",
            "CVR catalog: let 20,45% > sell 19,32% (+1,1pp) — ngược 1010.\n"
            "Snapshot: cho thuê ~19% contact-day vs bán ~8,55% trên ~7,5M listing-days.\n"
            "→ KPI & CRM không gộp “Nhà ở”.",
        ),
        ("Kênh &\nSeller", "Funnel &\nKênh lead"),
        (
            "Agent 76,1% supply · PR 54,0% vs Private 51,0% — khác nhà ở; cần playbook riêng cho B2B.\n"
            "Kênh: other_interaction 92,2% · view_phone 6,4% · chat 1,1% — hành vi so sánh trước khi liên hệ.",
            "8,38% session explicit contact · 50,8% events other_interaction (96% ad_view).\n"
            "76% explicit = view_phone (phone-first) · high-intent 5,98% vs broker/spam 12,38%.",
        ),
        ("C1 AREA", "C1 KPI"),
        (
            "30–50 m² thiếu cung.\n<30 m² +6,2pp.\n>200 m² thừa −12,6pp.",
            "8,38% session\nexplicit only.\n50,8% ad_view.",
        ),
        ("C2 SCHEMA", "C2 GEM"),
        ("NULL cao", "6.187"),
        (
            "ward/direction/\nproject_id thưa.\nlegal ~4,7% NULL.",
            "271% contact/PV\nexposure = 2.\nBoost priority.",
        ),
        ("C3 AGE", "C3 POCKET"),
        ("↓14 ngày", "310k"),
        (
            "Contacts/views\ngiảm sau 2 tuần.\nCần freshness boost.",
            "sell 30–50m²\nCVR 17,5%.\nDrag lớn nhất.",
        ),
        ("C4 PRICE", "C4 DEMOTE"),
        ("Gap", "4.295"),
        (
            "Bucket giá\nlệch cung/cầu\n(theo segment).",
            "Oversaturated\nexp. median 13.\nCap impression.",
        ),
        ("C5 GEO", "C5 BULK"),
        ("HCM+HN", "102/"),
        (
            "Tier-2 CVR cao\nsupply thấp.\nMở rộng địa lý.",
            "0,29 PV/listing\nseller spam.\nCap inventory.",
        ),
        ("C6 LET/SELL", "C6 UX"),
        ("94/6", "76%"),
        (
            "Hai thị trường\nmột feed — Sell\nbị chôn vùi.",
            "Phone-first +\ndeep-compare 60%.\nJourney mua dài.",
        ),
        (
            "Root cause: thiếu đúng phân khúc (30–100 m²) · schema thương mại yếu · Let/Sell chung feed · không tách B2B lifecycle.",
            "Root cause: feed ưu tiên visibility (ad_view) hơn verified converters · gem bị chôn · sell contact-day thấp · bulk seller loãng feed.",
        ),
        ("IMPACTS & ROADMAP", "STRATEGIES & IMPACTS"),
        ("2 / 4", "1 / 2"),
    ]
    _replace_in_slide(sc, reps)
    _set_shape_text(sc, 13, "19% /")
    _set_shape_text(sc, 14, "8,6%\ncontact-day\nlet / sell")
    _set_shape_text(sc, 45, "8,38%")
    _set_shape_text(
        sc,
        40,
        "8,38% session explicit · 50,8% events other_interaction (96% ad_view).\n"
        "76% explicit = view_phone (phone-first) · high-intent 5,98% vs broker/spam 12,38%.",
    )
    _set_shape_text(sc, 83, "STRATEGIES & IMPACTS")


def _apply_slide2(sc) -> None:
    """STRATEGIES & IMPACTS — layout from 1030 slide 4."""
    reps = [
        (
            "IMPACTS & ROADMAP — CATEGORY 1030 (VĂN PHÒNG / MẶT BẰNG)",
            "STRATEGIES & IMPACTS — CATEGORY 1020 (NHÀ Ở)",
        ),
        (
            "Mục tiêu: +25% CVR @ W4 · +45% @ M3 · +80% @ M12",
            "Mục tiêu: boost gem 6,2k · demote 4,3k · fix sell 30–50m² | Focus: sell HCM",
        ),
        ("+25%", "+93k"),
        ("CVR lift\nTuần 1–4", "contact events\n(boost gem)"),
        ("+45%", "2,7×"),
        ("CVR lift\nTháng 2–3", "vs 1010 boost\nopportunity"),
        ("+80%", "+9,7k"),
        ("CVR lift\nTháng 4–12", "sessions w/\ncontact (+1pp)"),
        ("+42%", "~322k"),
        ("Tier-2 geo\nuplift model", "imp-equiv saved\n(demote)"),
        ("+15%", "~6,2k"),
        ("Freshness\n0–14 ngày", "listings +\ncontact (pocket)"),
        ("Fill", "4,9k"),
        ("legal_status\n~4,7% NULL", "sell HCM\nin boost seg."),
        ("ROADMAP", "STRATEGIES"),
        ("Quick Wins (Tuần 1–4) → +25%", "Quick Wins (Tuần 1–2)"),
        (
            "① Bổ sung supply 50–100 m² (Let)\n② Hoàn thiện Legal Status (form + data fill)\n③ Image Quality Gate (≥5 ảnh)",
            "① Tách KPI layer — explicit_contact/session\n② Boost 6.187 HQ underexposed (ưu tiên sell HCM 4.853)\n③ Phone-first CTA — Gọi/Zalo sticky",
        ),
        ("Mid-term (Tháng 2–3) → +45%", "Mid-term (Tháng 2–4)"),
        (
            "④ Freshness Ranking (ưu tiên 0–14 ngày)\n⑤ Let vs Sell channel separation\n⑥ Tier-2 city listing expansion",
            "④ Demote 4.295 oversaturated — cap impression\n⑤ Completeness gate (house_type, floors, legal, width)\n⑥ Deep-compare UX (60% session)",
        ),
        ("Full Build (Tháng 4–12) → +80%", "Full (6 tháng)"),
        (
            "⑦ Legal Verification Trust Badge\n⑧ Sell Market Activation (HN + Đà Nẵng)\n⑨ Commercial Relevance Scoring v2 (ML)",
            "⑦ Android mobile-first (26% events)\n⑧ Cap bulk seller (~102 tin/seller)\n⑨ Fix pocket sell 30–50m² CVR 17,5→19,5%",
        ),
        (
            "1. Mid-size supply (30–100 m²): lấp gap area lớn nhất — CAO · Q1\n2. Tier-2 expansion: +42% CVR model → +XK revenue events — TB-CAO · Q2\n3. Freshness + legal fill + Sell activation — cumulative +80% @ M12",
            "1. Boost underexposed: +~92,8k contact events (exposure 2→4)\n2. Demote oversaturated: ~322k impression-equiv freed\n3. Session +1pp explicit: +~9,7k sessions · pocket +2pp: ~6,2k listings",
        ),
        (
            "commercial_relevance_1030 · ad_type_channel (let|sell) · area_price_segment · legal_trust_score · freshness_boost\nKhông copy playbook 1050/1040 — B2B cần tách kênh & segment theo diện tích×giá.",
            "health_segment_boost · ad_type_channel (let|sell) · sell_hcm_priority · phone_first_score · deep_compare_flag · seller_cap\nLever #1 recsys gem (~2,7× căn hộ); lever #2 sell 30–50m² + bulk seller.",
        ),
        (
            "Tuần 1: campaign seller 30–50 & 50–100 m² · bắt buộc legal_status + ≥5 ảnh · prototype tách Let/Sell feed.\nFocus tier-2: Đà Nẵng, Bình Dương — CVR cao, supply thấp (BCG Stars).",
            "Tuần 1–2: boost 6,2k gem + demote pilot · phone-first trên listing sell.\nFocus HCM sell (~70% supply, 4,9k tin trong segment boost). Caveat: explicit ≠ ad_view · sample 6–8%.",
        ),
        ("4 / 4", "2 / 2"),
    ]
    _replace_in_slide(sc, reps)
    _set_shape_text(sc, 1, "STRATEGIES & IMPACTS — CATEGORY 1020 (NHÀ Ở)")
    _set_shape_text(sc, 30, "Quick Wins (Tuần 1–2)")
    _set_shape_text(sc, 34, "Mid-term (Tháng 2–4)")
    _set_shape_text(sc, 38, "Full (6 tháng)")
    _set_shape_text(sc, 63, "STRATEGIES & IMPACTS")


def build() -> Path:
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Missing template: {TEMPLATE}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    shutil.copy2(TEMPLATE, OUT_PPTX)
    prs = Presentation(str(OUT_PPTX))

    # Remove overview (0) and strategies-only (2); keep S&C (1) and Impacts (3)
    _delete_slide(prs, 2)
    _delete_slide(prs, 0)

    _apply_slide1(prs.slides[0])
    _apply_slide2(prs.slides[1])

    prs.save(str(OUT_PPTX))
    return OUT_PPTX


def main() -> None:
    out = build()
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
