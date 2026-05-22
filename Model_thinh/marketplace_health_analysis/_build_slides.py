"""
Builds the condensed 5-slide deck — datathon_final_deck.pptx

Each slide is a dashboard combining multiple charts + content from the original 17-slide deck.
Font: Arial (universal — Mac/Win/Linux + full Vietnamese diacritic support).
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Paths ─────────────────────────────────────────────────────────────────────
ROOT     = Path('/Volumes/mavuong3112/Datathon_Data/marketplace_health_analysis')
CHARTS   = ROOT / 'slides' / 'charts'
OUT_PPTX = ROOT / 'slides' / 'datathon_final_deck.pptx'
OUT_PPTX.parent.mkdir(exist_ok=True, parents=True)

# ── Palette ───────────────────────────────────────────────────────────────────
SLATE      = RGBColor(0x1B, 0x25, 0x38)
SLATE_DARK = RGBColor(0x0F, 0x16, 0x24)
INDIGO     = RGBColor(0x3F, 0x51, 0xB5)
SKY        = RGBColor(0x00, 0xB0, 0xFF)
MINT       = RGBColor(0x00, 0xBF, 0xA5)
CORAL      = RGBColor(0xFF, 0x52, 0x52)
AMBER      = RGBColor(0xFF, 0xB7, 0x4D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x10, 0x20, 0x30)
MUTED      = RGBColor(0x60, 0x7D, 0x8B)
SOFT_GREY  = RGBColor(0xEC, 0xEF, 0xF3)
PALE_MINT  = RGBColor(0xE0, 0xF2, 0xF1)
PALE_AMBER = RGBColor(0xFF, 0xF3, 0xE0)
PALE_CORAL = RGBColor(0xFF, 0xEB, 0xEE)

FONT = 'Arial'
FONT_MONO = 'Courier New'

# ── 16:9 widescreen ───────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.500)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

SECTIONS = [
    'OVERVIEW & DATA',
    'ARCHITECTURE',
    'PERFORMANCE',
    'MARKETPLACE HEALTH',
    'PRODUCTION & FUTURE',
]
TOTAL = 5

# ═════════════════════════════════════════════════════════════════════════════
# Primitive helpers
# ═════════════════════════════════════════════════════════════════════════════
def rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def rounded_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=0.5, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def txt(slide, x, y, w, h, text, size=11, bold=False, italic=False,
        color=DARK_TEXT, align='left', valign='top', font=FONT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, Emu(0))
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP,
                          'middle': MSO_ANCHOR.MIDDLE,
                          'bottom': MSO_ANCHOR.BOTTOM}[valign]
    aln = {'left': PP_ALIGN.LEFT,
           'center': PP_ALIGN.CENTER,
           'right': PP_ALIGN.RIGHT}[align]
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = aln
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb

def pic(slide, x, y, w, h, path):
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)

def bullets(slide, x, y, w, h, items, size=10, color=DARK_TEXT,
            bullet_color=None, bullet='▸', line_spacing=1.30, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, Emu(0))
    bc = bullet_color or color
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r1 = p.add_run(); r1.text = f'{bullet}  '
        r1.font.name = font; r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = bc
        r2 = p.add_run(); r2.text = item
        r2.font.name = font; r2.font.size = Pt(size); r2.font.color.rgb = color
    return tb

# ═════════════════════════════════════════════════════════════════════════════
# Slide chrome (header, title, footer)
# ═════════════════════════════════════════════════════════════════════════════
def header_strip(slide, section_idx, page_num, total):
    rect(slide, 0, 0, SW, Inches(0.50), fill=SLATE)
    rect(slide, 0, Inches(0.50), SW, Inches(0.04), fill=CORAL)
    bread = f' DATATHON 2026   ›   {SECTIONS[section_idx]} '
    txt(slide, Inches(0.45), Inches(0.10),
        Inches(8.0), Inches(0.30), bread,
        size=10, bold=True, color=WHITE, valign='middle')
    txt(slide, SW - Inches(3.5), Inches(0.10),
        Inches(2.55), Inches(0.30),
        'CHỢ TỐT BĐS · The Gridbreakers',
        size=9, color=WHITE, valign='middle', align='right')
    rect(slide, SW - Inches(0.75), Inches(0.10),
         Inches(0.55), Inches(0.30), fill=CORAL)
    txt(slide, SW - Inches(0.75), Inches(0.10),
        Inches(0.55), Inches(0.30), f'{page_num:02d}',
        size=11, bold=True, color=WHITE, align='center', valign='middle')

def footer_progress(slide, page_num, total):
    bar_y = SH - Inches(0.18)
    rect(slide, Inches(0.45), bar_y, SW - Inches(0.90), Inches(0.06), fill=SOFT_GREY)
    frac = page_num / total
    rect(slide, Inches(0.45), bar_y,
         (SW - Inches(0.90)) * frac, Inches(0.06), fill=CORAL)
    txt(slide, Inches(0.45), bar_y - Inches(0.22),
        Inches(3.0), Inches(0.20),
        f'{page_num} of {total}',
        size=7.5, color=MUTED)

def slide_title(slide, title, subtitle=None):
    txt(slide, Inches(0.45), Inches(0.72),
        Inches(12.5), Inches(0.55), title,
        size=22, bold=True, color=SLATE_DARK, line_spacing=1.0)
    if subtitle:
        rect(slide, Inches(0.45), Inches(1.22),
             Inches(0.40), Inches(0.06), fill=CORAL)
        txt(slide, Inches(0.95), Inches(1.13),
            Inches(11.5), Inches(0.26), subtitle,
            size=10.5, italic=True, color=MUTED)

def base_slide(section_idx, page_num, total, title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    header_strip(s, section_idx, page_num, total)
    slide_title(s, title, subtitle)
    footer_progress(s, page_num, total)
    return s

# ═════════════════════════════════════════════════════════════════════════════
# Components
# ═════════════════════════════════════════════════════════════════════════════
def kpi_tile(slide, x, y, w, h, value, label, sublabel=None, color=INDIGO,
             value_size=42, label_size=10, sub_size=8.5):
    rounded_rect(slide, x, y, w, h, fill=color)
    txt(slide, x + Inches(0.12), y + Inches(0.10),
        w - Inches(0.24), h * 0.52, value,
        size=value_size, bold=True, color=WHITE, align='center', valign='middle',
        line_spacing=1.0)
    txt(slide, x + Inches(0.10), y + h*0.58,
        w - Inches(0.20), Inches(0.24), label,
        size=label_size, bold=True, color=WHITE, align='center')
    if sublabel:
        txt(slide, x + Inches(0.10), y + h*0.78,
            w - Inches(0.20), h*0.20, sublabel,
            size=sub_size, color=WHITE, align='center', italic=True,
            line_spacing=1.10)

def panel(slide, x, y, w, h, title, items=None, body_text=None,
          color=INDIGO, body_size=9.5, mono=False, line_spacing=1.30,
          header_h=0.32):
    rounded_rect(slide, x, y, w, h, fill=WHITE, line=SOFT_GREY, line_w=0.75)
    hb = Inches(header_h)
    rounded_rect(slide, x, y, w, hb, fill=color)
    txt(slide, x + Inches(0.15), y, w - Inches(0.30), hb, title,
        size=10, bold=True, color=WHITE, valign='middle')
    bx, by = x + Inches(0.15), y + hb + Inches(0.05)
    bw, bh = w - Inches(0.30), h - hb - Inches(0.10)
    if items is not None:
        bullets(slide, bx, by, bw, bh, items,
                size=body_size, color=DARK_TEXT, bullet_color=color,
                line_spacing=line_spacing)
    elif body_text is not None:
        txt(slide, bx, by, bw, bh, body_text,
            size=body_size, color=DARK_TEXT,
            font=FONT_MONO if mono else FONT,
            line_spacing=line_spacing)

def callout(slide, x, y, w, h, text_str, label='KEY INSIGHT',
            color=AMBER, bg=PALE_AMBER, body_size=9.5):
    rounded_rect(slide, x, y, w, h, fill=bg, line=color, line_w=1.2)
    rect(slide, x, y, Inches(0.10), h, fill=color)
    txt(slide, x + Inches(0.22), y + Inches(0.06),
        w - Inches(0.32), Inches(0.22), label,
        size=8.5, bold=True, color=color)
    txt(slide, x + Inches(0.22), y + Inches(0.28),
        w - Inches(0.32), h - Inches(0.36), text_str,
        size=body_size, color=DARK_TEXT, line_spacing=1.30)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — OVERVIEW + PROBLEM + DATA  (covers original slides 1, 2, 3)
# ═════════════════════════════════════════════════════════════════════════════
def s1_overview():
    s = base_slide(0, 1, TOTAL,
                   'Tổng quan · Bài toán · Dataset',
                   '161,568 users · 3.1M items · Recall@10 = 0.2441 (Public LB) · 42× lift vs popularity')

    # Row 1 — Hero + 3 KPIs + Pipeline mini
    hero_y = Inches(1.55); hero_h = Inches(1.80)

    # Hero KPI (left)
    kpi_tile(s, Inches(0.45), hero_y, Inches(3.10), hero_h,
             '0.2441', 'RECALL@10 · PUBLIC LB',
             'v15 stage15 · best score',
             color=SLATE, value_size=48, label_size=10)

    # 3 supporting compact KPIs (middle)
    sub_kpis = [
        ('42×',     'LIFT vs POPULARITY',    'baseline 0.0058', MINT),
        ('161K',    'TEST USERS',            'top-10 cần predict', SKY),
        ('3.1M',    'ITEMS POOL',            '5 categories', INDIGO),
    ]
    skx, sky_, skw, skh = Inches(3.70), hero_y, Inches(2.50), hero_h
    gap = Inches(0.08)
    for i, (v, l, sub, col) in enumerate(sub_kpis):
        x = skx + (skw + gap) * 0
        # 3 stacked vertically would be too tight — use horizontal layout
    # Switch to horizontal for sub_kpis
    skw = Inches(2.55); skh = Inches(0.58)
    skx = Inches(3.70)
    for i, (v, l, sub, col) in enumerate(sub_kpis):
        y = hero_y + (skh + Inches(0.05)) * i
        rounded_rect(s, skx, y, Inches(7.78), skh, fill=col)
        txt(s, skx + Inches(0.15), y + Inches(0.05),
            Inches(1.5), Inches(0.50), v,
            size=22, bold=True, color=WHITE, valign='middle')
        txt(s, skx + Inches(1.80), y + Inches(0.04),
            Inches(2.50), Inches(0.26), l,
            size=10, bold=True, color=WHITE)
        txt(s, skx + Inches(1.80), y + Inches(0.30),
            Inches(2.50), Inches(0.26), sub,
            size=8.5, color=WHITE, italic=True)
        # Pipeline summary tucked right of subkpi row 0
    # Pipeline mini panel (right of hero+sub_kpis)
    panel(s, Inches(11.65), hero_y, Inches(1.28), hero_h,
          'PIPELINE', color=INDIGO,
          items=['5 retrievers',
                 'ALS · EASE · ItemCF',
                 'SASRec · CBF',
                 '54 features',
                 'LGBM × 10 seeds',
                 '+ XGBoost blend'],
          body_size=8, line_spacing=1.30, header_h=0.28)
    # Actually pipeline panel should be wider. Let me redo layout.

# Redo s1_overview with a cleaner layout
def s1_overview():
    s = base_slide(0, 1, TOTAL,
                   'Tổng quan · Bài toán · Dataset',
                   '161,568 users · 3.1M items · Recall@10 = 0.2441 (Public LB) · 42× lift vs popularity')

    # Row 1 (y=1.55-3.40): Hero + 3 sub-KPIs + Pipeline panel
    r1_y = Inches(1.55); r1_h = Inches(1.85)

    # Hero (left)
    kpi_tile(s, Inches(0.45), r1_y, Inches(3.20), r1_h,
             '0.2441', 'RECALL@10 · PUBLIC LB',
             'v15 stage15 · best score',
             color=SLATE, value_size=48, label_size=10, sub_size=8.5)

    # 3 sub-KPIs horizontally between hero and pipeline panel
    sub = [
        ('42×',   'LIFT',         '↑ vs popularity',  MINT),
        ('161K',  'TEST USERS',   'top-10 / user',    SKY),
        ('3.1M',  'ITEMS',        '5 categories',     INDIGO),
    ]
    sx0 = Inches(3.85); skw = Inches(2.05); skh = r1_h; gap = Inches(0.08)
    for i, (v, l, sb, col) in enumerate(sub):
        x = sx0 + (skw + gap) * i
        kpi_tile(s, x, r1_y, skw, skh, v, l, sb,
                 color=col, value_size=32, label_size=9.5, sub_size=8)

    # Pipeline panel (right)
    panel(s, Inches(10.27), r1_y, Inches(2.66), r1_h,
          'PIPELINE OVERVIEW', color=INDIGO,
          items=[
              '5 retrievers (top-100 each)',
              'ALS · EASE · ItemCF · SASRec · CBF',
              '54 features (4 groups)',
              'LGBM LambdaRank ×10 seeds',
              '+ XGBoost blend (0.65/0.35)',
              'Freshness + cat boost',
              'Cold-user fallback',
          ], body_size=8.5, line_spacing=1.32, header_h=0.30)

    # Row 2 (y=3.55-5.00): Problem + Data
    r2_y = Inches(3.55); r2_h = Inches(1.50)

    # Problem formal equation card (left)
    rounded_rect(s, Inches(0.45), r2_y, Inches(6.50), r2_h,
                 fill=RGBColor(0xE3, 0xF2, 0xFD), line=SKY, line_w=1.0)
    txt(s, Inches(0.60), r2_y + Inches(0.10),
        Inches(6.20), Inches(0.25),
        '◆ FORMAL OBJECTIVE',
        size=9.5, bold=True, color=SKY)
    txt(s, Inches(0.60), r2_y + Inches(0.40),
        Inches(6.20), Inches(0.45),
        'argmax_θ    E_u [ |f_θ(u) ∩ G_u| / |G_u| ]',
        size=16, bold=True, color=SLATE_DARK, font=FONT_MONO, align='center')
    txt(s, Inches(0.60), r2_y + Inches(0.95),
        Inches(6.20), Inches(0.50),
        'f_θ : U → I^10   ·   G_u = ground truth (10/04 → 07/05, private)\n'
        'Metric chính: Recall@10  ·  Tie-break: NDCG@10',
        size=9, color=DARK_TEXT, align='center', line_spacing=1.30)

    # 4 mini data KPIs (right of problem)
    mk = [
        ('U',  '161K',     SKY),
        ('I',  '3.1M',     INDIGO),
        ('T',  '152d',     MINT),
        ('Δt', '28d',      CORAL),
    ]
    mkx = Inches(7.05); mkw = Inches(1.43); mkh = r2_h; mkgap = Inches(0.08)
    for i, (sym, val, col) in enumerate(mk):
        x = mkx + (mkw + mkgap) * i
        rounded_rect(s, x, r2_y, mkw, mkh, fill=WHITE, line=col, line_w=1.2)
        rect(s, x, r2_y, mkw, Inches(0.35), fill=col)
        txt(s, x, r2_y, mkw, Inches(0.35), sym,
            size=14, bold=True, color=WHITE, italic=True,
            align='center', valign='middle')
        txt(s, x, r2_y + Inches(0.40),
            mkw, Inches(0.55), val,
            size=20, bold=True, color=col, align='center', valign='middle')
        labels_map = {'U':'users', 'I':'items', 'T':'train', 'Δt':'eval'}
        txt(s, x, r2_y + Inches(1.00),
            mkw, Inches(0.40), labels_map[sym],
            size=8.5, color=DARK_TEXT, italic=True, align='center')

    # Row 3 (y=5.15-6.95): 4 detail panels
    r3_y = Inches(5.15); r3_h = Inches(1.85); r3w = Inches(3.07); r3gap = Inches(0.08)
    panels_data = [
        ('POSITIVE EVENTS', MINT, [
            'view_phone',
            'contact_chat',
            'contact_zalo',
            'contact_sms',
            'other_interaction',
        ]),
        ('5 CATEGORIES', INDIGO, [
            '1010 — Phòng trọ',
            '1020 — Căn hộ',
            '1030 — Nhà ở',
            '1040 — Đất nền',
            '1050 — Dự án mới',
        ]),
        ('CONSTRAINTS', CORAL, [
            'Không leak (< 10/04)',
            'Không data ngoài',
            '5 subs/đội/day UTC',
            'Top-10 phân biệt',
            'Random seed cố định',
        ]),
        ('KEY INSIGHTS', AMBER, [
            'Organic fairness 3.3×',
            'private 16.5% pool',
            '→ 54.9% exposed',
            'Freshness 1.1% (gap)',
            'HHI < 100 (diverse)',
        ]),
    ]
    for i, (title, col, items) in enumerate(panels_data):
        x = Inches(0.45) + (r3w + r3gap) * i
        panel(s, x, r3_y, r3w, r3_h, title, color=col,
              items=items, body_size=9, line_spacing=1.40, header_h=0.30)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ARCHITECTURE + MATH + FEATURES (covers original 4, 5, 6)
# ═════════════════════════════════════════════════════════════════════════════
def s2_architecture():
    s = base_slide(1, 2, TOTAL,
                   'Solution Architecture · Math · Features',
                   'Two-stage paradigm · 5 retrievers · 54 features · LGBM + XGBoost ensemble')

    # Row 1: 3 stage cards horizontally
    sy = Inches(1.55); sh = Inches(2.15); sw = Inches(4.15); gap = Inches(0.08)
    stages = [
        ('STAGE 1', 'CANDIDATE GENERATION', SKY, [
            '5 retrievers · top-100 each',
            'ALS (factors=512, α=40)',
            'EASE (closed-form, λ=200)',
            'ItemCF (time-decayed cosine)',
            'SASRec (self-attention seq)',
            'CBF (PhoBERT 768d)',
            '→ merge & cap 700/user',
        ]),
        ('STAGE 2', 'LEARNING TO RANK', INDIGO, [
            '54 features → ranker',
            'LGBM LambdaRank · depth 8',
            'num_leaves 95 · lr 0.05',
            'n_estimators 500',
            '10-seed ensemble (variance ↓)',
            'XGBoost (depth 6, 1000 rounds)',
            '→ ranked_predictions',
        ]),
        ('STAGE 3', 'BLEND + POST-PROCESS', MINT, [
            'Per-user min-max norm',
            's_final = 0.65·LGBM + 0.35·XGB',
            '+ β_f·1[fresh ≤7d]',
            '+ β_c·1[cat=pref]',
            'Cold-user fallback (cat-weighted)',
            '→ top-10 submission.csv',
        ]),
    ]
    for i, (num, title, col, items) in enumerate(stages):
        x = Inches(0.45) + (sw + gap) * i
        rounded_rect(s, x, sy, sw, sh, fill=WHITE, line=col, line_w=1.5)
        rect(s, x, sy, sw, Inches(0.55), fill=col)
        txt(s, x + Inches(0.18), sy + Inches(0.06),
            Inches(1.0), Inches(0.20), num,
            size=8.5, bold=True, color=WHITE)
        txt(s, x + Inches(0.18), sy + Inches(0.22),
            sw - Inches(0.36), Inches(0.30), title,
            size=12, bold=True, color=WHITE, valign='middle')
        bullets(s, x + Inches(0.22), sy + Inches(0.68),
                sw - Inches(0.40), sh - Inches(0.78),
                items, size=9, color=DARK_TEXT,
                bullet_color=col, line_spacing=1.35)
    # Arrows
    arr_y = sy + sh/2 - Inches(0.10)
    for i in range(2):
        ax = Inches(0.45) + sw * (i+1) + gap * i
        tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, arr_y, gap, Inches(0.20))
        tri.fill.solid(); tri.fill.fore_color.rgb = MUTED
        tri.line.fill.background(); tri.shadow.inherit = False

    # Row 2: Math (left) + Feature importance chart (middle) + Taxonomy (right)
    ry = Inches(3.90); rh = Inches(3.10)

    # Math panel
    panel(s, Inches(0.45), ry, Inches(5.00), rh,
          'MATH NOTATION · 5 RETRIEVERS + RANKER', color=SLATE,
          body_size=8.5, mono=True,
          body_text=(
              'ALS [1]      L = Σ c_ui (p_ui - u_u^T v_i)^2 + λ(||U||^2 + ||V||^2)\n'
              '             c_ui = 1 + α·r_ui ;  α=40 ; factors=512\n\n'
              'EASE [2]     B = (G + λI)^-1 ;  B_ii <- 0 ;  G = X^T X ;  λ=200\n\n'
              'ItemCF [3]   s(i,j) = Σ_u w(t)·1[u,i]·1[u,j] / sqrt(...)\n'
              '             w(t) = exp(-β·Δt) ;  β=0.005 (~140d half-life)\n\n'
              'SASRec [4]   SA(Q,K,V) = softmax(QK^T / sqrt(d))·V\n'
              '             d=64, heads=2, blocks=2, L_max=50\n\n'
              'CBF [8]      e_i = PhoBERT(title_i) in R^768 ;  cos(e_u, e_i)\n\n'
              'LambdaRank [5,6]\n'
              '   L = Σ log(1 + exp(-σ(s_i - s_j))) · |Δ NDCG_ij|\n\n'
              's_final = 0.65·s_LGBM_norm + 0.35·s_XGB_norm\n'
              '          + β_f·1[fresh] + β_c·1[cat_match]'
          ), line_spacing=1.22)

    # Feature importance chart (middle)
    pic(s, Inches(5.60), ry, Inches(4.00), rh,
        CHARTS / 'chart_feature_importance.png')

    # Feature taxonomy (right)
    panel(s, Inches(9.75), ry, Inches(3.20), rh,
          'FEATURE TAXONOMY · 54', color=AMBER,
          items=[
              'USER-LEVEL · 15',
              'n_pos_events, intent_score_log, pref_category',
              '',
              'ITEM-LEVEL · 18',
              'trend_pos, days_since_post, item_quality',
              '',
              'CROSS user×item · 12',
              'is_repeat, ui_match_category, ui_district',
              '',
              'SEQUENCE/INTENT · 9',
              'last_click_days, dwell_avg, first_click_match',
          ], body_size=8.5, line_spacing=1.20, header_h=0.30)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PERFORMANCE DASHBOARD (covers original 7, 8, 9)
# ═════════════════════════════════════════════════════════════════════════════
def s3_performance():
    s = base_slide(2, 3, TOTAL,
                   'Performance Dashboard · Recall@10 = 0.2441',
                   'LB progression · Personalization 42× · Score diagnostics · Top features')

    # Top — 4 KPI tiles (compact row)
    kx0 = Inches(0.45); kw = Inches(3.07); kh = Inches(1.30); kgap = Inches(0.08)
    kpis = [
        ('0.2441',  'RECALL@10 (LB)',      'v15 stage15 best', SLATE),
        ('42×',     'PERSONALIZATION LIFT', 'vs popularity 0.0058', MINT),
        ('+11.8%',  'GAIN vs v1',           '0.2184 → 0.2441 in 11d', INDIGO),
        ('1.2/10',  'AVG OVERLAP',          'với popularity top-10', AMBER),
    ]
    for i, (v, l, sb, col) in enumerate(kpis):
        x = kx0 + (kw + kgap) * i
        kpi_tile(s, x, Inches(1.55), kw, kh,
                 v, l, sb, color=col, value_size=28, label_size=9.5, sub_size=8)

    # Middle — Performance combo dashboard chart (4-panel)
    pic(s, Inches(0.45), Inches(2.95),
        Inches(9.50), Inches(3.80),
        CHARTS / 'dashboard_performance.png')

    # Right column — Insights stack
    rx = Inches(10.10); rw = Inches(2.83)
    ry = Inches(2.95); rh = Inches(1.20)
    cgap = Inches(0.05)
    insights = [
        ('STRATEGIC DECISION', AMBER,
         'Diminishing returns ở 0.244 → dừng tối ưu accuracy, chuyển focus sang Marketplace Health.'),
        ('PERSONALIZATION PROOF', MINT,
         '80-90% slot trong top-10 là personalized, không phải lazy popularity.'),
        ('SCORE STRUCTURE', SKY,
         'Top-3 confident · Mid 4-10 fragile · Tail = inherent uncertainty.'),
    ]
    for i, (lbl, col, body) in enumerate(insights):
        y = ry + (rh + cgap) * i
        callout(s, rx, y, rw, rh, body, label=lbl, color=col,
                bg=WHITE, body_size=8.5)

    # Bottom — Honest disclaimer banner
    db_y = Inches(6.85 - 0.05)
    rounded_rect(s, Inches(0.45), db_y, Inches(12.48), Inches(0.42),
                 fill=PALE_AMBER, line=AMBER, line_w=1.0)
    rect(s, Inches(0.45), db_y, Inches(0.08), Inches(0.42), fill=AMBER)
    txt(s, Inches(0.65), db_y, Inches(12.28), Inches(0.42),
        '⚠  HẠN CHẾ TRUNG THỰC:  NDCG@10 nội bộ saturate về 1.0 (feature leakage) · Recall per segment không đo được (GT private) · '
        'Variant A/B Recall = proxy qua overlap · Public LB chỉ 5% GT',
        size=8.5, italic=True, color=DARK_TEXT, valign='middle')

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MARKETPLACE HEALTH DASHBOARD (covers original 10-14) ⭐
# ═════════════════════════════════════════════════════════════════════════════
def s4_health():
    s = base_slide(3, 4, TOTAL,
                   'Marketplace Health Dashboard · 5 trục đánh giá',
                   'Scorecard · Trade-off · Concentration · Age skew · Fairness — Organic 3.3× private boost')

    # Left — combo health dashboard (4-panel)
    pic(s, Inches(0.45), Inches(1.55),
        Inches(8.40), Inches(4.50),
        CHARTS / 'dashboard_health.png')

    # Right top — Fairness donut highlight
    pic(s, Inches(8.95), Inches(1.55),
        Inches(4.00), Inches(2.20),
        CHARTS / 'chart_seller_donut.png')

    # Right middle — Fairness key tile
    rounded_rect(s, Inches(8.95), Inches(3.85),
                 Inches(4.00), Inches(0.78), fill=MINT)
    txt(s, Inches(9.10), Inches(3.90),
        Inches(3.70), Inches(0.25), 'ORGANIC FAIRNESS',
        size=9, bold=True, color=WHITE)
    txt(s, Inches(9.10), Inches(4.15),
        Inches(3.70), Inches(0.45),
        'private 16.5% pool → 54.9% exposed  ·  3.3× boost',
        size=11, bold=True, color=WHITE)

    # Right bottom — Verdict callouts
    cb_y = Inches(4.75); cb_h = Inches(1.30)
    rounded_rect(s, Inches(8.95), cb_y, Inches(4.00), cb_h,
                 fill=PALE_MINT, line=MINT, line_w=1.2)
    rect(s, Inches(8.95), cb_y, Inches(0.10), cb_h, fill=MINT)
    txt(s, Inches(9.15), cb_y + Inches(0.08),
        Inches(3.75), Inches(0.22),
        '💡  WHY private over-represented?',
        size=9, bold=True, color=MINT)
    bullets(s, Inches(9.15), cb_y + Inches(0.32),
            Inches(3.75), cb_h - Inches(0.40),
            ['Contact rate cao hơn (no commission)',
             'Title/desc match user intent tốt hơn',
             'Pricing thực tế hơn (no agent markup)',
             '→ Không cần fairness constraint thêm'],
            size=8.5, color=DARK_TEXT, bullet_color=MINT, line_spacing=1.30)

    # Bottom — Trade-off matrix table
    tm_y = Inches(6.15)
    rounded_rect(s, Inches(0.45), tm_y, Inches(12.48), Inches(1.13),
                 fill=WHITE, line=SLATE, line_w=0.75)
    rect(s, Inches(0.45), tm_y, Inches(12.48), Inches(0.28), fill=SLATE)
    txt(s, Inches(0.65), tm_y, Inches(12.28), Inches(0.28),
        '◆ TRADE-OFF MATRIX  ·  Decision per scenario',
        size=10, bold=True, color=WHITE, valign='middle')
    txt(s, Inches(0.65), tm_y + Inches(0.35),
        Inches(12.20), Inches(0.75),
        'STRATEGY                              Δ Recall (est)    Δ Freshness    Δ Seller Cov    DECISION\n'
        'Baseline (v15)                              —               —              —          Current LB best\n'
        '+ Freshness boost +5%                   −0.7% (0.2370)   +180% (×2.7)    +7%          ✓ YES production\n'
        '+ Hard seller cap ≤2                    −56%  (0.1077)   −18%            +7%          ✗ TOO AGGRESSIVE\n'
        '+ Soft seller cap = 3 (proposed)        −5–10% (est)     ~0              +3%          ⟳ NEEDS A/B TEST',
        size=7.8, color=DARK_TEXT, font=FONT_MONO, line_spacing=1.30)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PRODUCTION + ROADMAP + LIMITATIONS (covers original 15, 16, 17)
# ═════════════════════════════════════════════════════════════════════════════
def s5_production():
    s = base_slide(4, 5, TOTAL,
                   'Production Readiness · Roadmap · Honest Limitations',
                   'Serving · Latency · Failure modes · Roadmap 30/60/90 · Self-assessment 4 trục BGK')

    # Row 1 — 4 production KPI tiles
    kx0 = Inches(0.45); kw = Inches(3.07); kh = Inches(1.10); kgap = Inches(0.08)
    kpis = [
        ('< 5 ms',   'INFERENCE LATENCY',   'LGBM per user batch', MINT),
        ('< 20 MB',  'MODEL FOOTPRINT',     'fit 1 pod · no GPU', INDIGO),
        ('< 1 min',  'THROUGHPUT (8 cores)', 'full 161K users', SKY),
        ('~$0.20',   'RETRAIN COST/WEEK',   '3-4h AWS spot', AMBER),
    ]
    for i, (v, l, sb, col) in enumerate(kpis):
        x = kx0 + (kw + kgap) * i
        kpi_tile(s, x, Inches(1.55), kw, kh,
                 v, l, sb, color=col, value_size=22, label_size=9, sub_size=7.5)

    # Row 2 — Architecture (left) + Failure modes (right)
    ry = Inches(2.80); rh = Inches(2.45)

    panel(s, Inches(0.45), ry, Inches(6.30), rh,
          'SERVING ARCHITECTURE · 2-TIER', color=SLATE,
          body_size=8, mono=True,
          body_text=(
              'REAL-TIME PATH  (< 100 ms p95)\n'
              '  Feature Store (Redis/Feast)  ↔  10-20 ms\n'
              '  ALS + ItemCF retrieval         ~30-50 ms (top-200)\n'
              '  LGBM ranker in-memory          < 5 ms\n'
              '  → top-10 returned\n\n'
              'OFFLINE PATH  (batch)\n'
              '  trending_pop          ── hourly\n'
              '  cold_top10            ── daily\n'
              '  ALS/EASE/SASRec       ── weekly\n'
              '  LGBM ranker           ── weekly\n'
              '  Full feature rebuild  ── monthly'
          ), line_spacing=1.25, header_h=0.30)

    panel(s, Inches(6.85), ry, Inches(6.08), rh,
          'FAILURE MODES → GRACEFUL DEGRADATION', color=CORAL,
          body_size=8, mono=True,
          body_text=(
              'Candidate gen down    → trending_pop fallback\n'
              'LGBM corrupt         → blend_score thuần (−30% R)\n'
              'Feature store timeout → cold_top10 (cached pkl)\n'
              'posted_date null      → days_since=999 default\n'
              'Cold user no profile  → category-weighted pool\n'
              'Traffic spike 3×      → auto-scale + cached top-10\n'
              'Train/serve skew      → daily diff alert → retrain\n\n'
              'TRIẾT LÝ:  KHÔNG bao giờ trả top-10 rỗng.\n'
              'Worst case → cold_top10.pkl (always cached).'
          ), line_spacing=1.30, header_h=0.30)

    # Row 3 — Roadmap 30/60/90 (left) + Self-assessment (middle) + Limitations (right)
    ry2 = Inches(5.35); rh2 = Inches(1.93)

    # Roadmap mini
    rmx = Inches(0.45); rmw = Inches(5.40)
    rounded_rect(s, rmx, ry2, rmw, rh2, fill=WHITE, line=SOFT_GREY, line_w=0.75)
    rounded_rect(s, rmx, ry2, rmw, Inches(0.32), fill=SLATE)
    txt(s, rmx + Inches(0.15), ry2, rmw - Inches(0.30), Inches(0.32),
        '◆ ROADMAP · 30 / 60 / 90 DAYS', size=10, bold=True, color=WHITE,
        valign='middle')
    # 3 sub-columns
    sub_w = Inches(1.66); sub_x0 = rmx + Inches(0.12)
    sub_gap = Inches(0.10)
    phases = [
        ('30d', 'CANARY DEPLOY', SKY, ['Variant A canary', 'Monitoring dashboard', 'A/B test infra']),
        ('60d', 'TWO-TOWER',     INDIGO, ['user×item tower PoC', 'Online learning 24h', 'Seller cold-start']),
        ('90d', 'GNN + MULTI-OBJ', MINT, ['GNN graph learning', 'Multi-obj optimization', 'Cross-category']),
    ]
    for i, (num, t, col, items) in enumerate(phases):
        x = sub_x0 + (sub_w + sub_gap) * i
        rounded_rect(s, x, ry2 + Inches(0.40), sub_w, rh2 - Inches(0.48),
                     fill=WHITE, line=col, line_w=1.0)
        rect(s, x, ry2 + Inches(0.40), sub_w, Inches(0.36), fill=col)
        txt(s, x + Inches(0.10), ry2 + Inches(0.40),
            Inches(0.5), Inches(0.36), num,
            size=12, bold=True, color=WHITE, valign='middle')
        txt(s, x + Inches(0.55), ry2 + Inches(0.40),
            sub_w - Inches(0.65), Inches(0.36), t,
            size=8, bold=True, color=WHITE, valign='middle')
        bullets(s, x + Inches(0.12), ry2 + Inches(0.82),
                sub_w - Inches(0.24), rh2 - Inches(0.92),
                items, size=7.5, color=DARK_TEXT, bullet_color=col,
                line_spacing=1.30)

    # Self-assessment (middle)
    sax = Inches(5.95); saw = Inches(3.65)
    rounded_rect(s, sax, ry2, saw, rh2, fill=WHITE, line=SOFT_GREY, line_w=0.75)
    rounded_rect(s, sax, ry2, saw, Inches(0.32), fill=SLATE)
    txt(s, sax + Inches(0.15), ry2, saw - Inches(0.30), Inches(0.32),
        '◆ SELF-ASSESSMENT · 4 TRỤC BGK', size=10, bold=True, color=WHITE,
        valign='middle')
    crit = [
        ('Thiết kế',       4.0, SKY),
        ('Hiệu suất',      4.0, INDIGO),
        ('Health',         4.5, MINT),
        ('Production',     4.0, AMBER),
    ]
    bar_y0 = ry2 + Inches(0.42)
    bar_h = Inches(0.30)
    for i, (lbl, score, col) in enumerate(crit):
        y = bar_y0 + i * (bar_h + Inches(0.06))
        txt(s, sax + Inches(0.18), y, Inches(1.10), bar_h, lbl,
            size=9, bold=True, color=DARK_TEXT, valign='middle')
        bw_max = Inches(1.65)
        rect(s, sax + Inches(1.35), y + Inches(0.08),
             bw_max, Inches(0.14), fill=SOFT_GREY)
        rect(s, sax + Inches(1.35), y + Inches(0.08),
             bw_max * (score/5.0), Inches(0.14), fill=col)
        txt(s, sax + Inches(3.05), y, Inches(0.55), bar_h, f'{score}/5',
            size=9, bold=True, color=col, valign='middle', align='right')

    # Limitations (right)
    panel(s, Inches(9.75), ry2, Inches(3.18), rh2,
          'HẠN CHẾ TRUNG THỰC', color=CORAL,
          items=[
              'NDCG offline saturate (leak)',
              'Recall/segment không đo được',
              'Variant A/B = overlap proxy',
              'Public LB chỉ 5% GT',
              'Latency = local benchmark',
          ], body_size=8.5, line_spacing=1.35, header_h=0.30)

# ─────────────────────────────────────────────────────────────────────────────
# BUILD
# ─────────────────────────────────────────────────────────────────────────────
slide_funcs = [
    s1_overview,
    s2_architecture,
    s3_performance,
    s4_health,
    s5_production,
]

for i, fn in enumerate(slide_funcs, 1):
    fn()
    print(f'  ✓ slide {i:02d}/{TOTAL}  {fn.__name__}')

prs.save(str(OUT_PPTX))
size_kb = OUT_PPTX.stat().st_size // 1024
print(f'\n✓ Saved: {OUT_PPTX}  ({size_kb} KB · {len(slide_funcs)} slides)')
