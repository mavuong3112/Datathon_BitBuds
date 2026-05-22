#!/usr/bin/env python3
"""
Build Macro/Demand PPTX from CSV seeds + existing notebook exports.

Data sources (no parquet scan unless seed CSV missing):
  - eda_weekly_stl_tet_regime  → outputs/macro_demand_deck/csv/stl_*.csv
  - eda_demand_category_device → outputs/eda_demand_category_device/02_*.csv
  - eda_listing_snapshot_cross_pie → outputs/eda_listing_snapshot/_cross/01_*.csv
"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns

DATA_ROOT = Path(__file__).resolve().parent
CSV_DIR = DATA_ROOT / "outputs" / "macro_demand_deck" / "csv"
CHART_DIR = DATA_ROOT / "outputs" / "macro_demand_deck" / "charts"
PPTX_PATH = DATA_ROOT / "outputs" / "macro_demand_deck" / "Macro_Demand_EDA_Deck.pptx"

CSV_DAILY = CSV_DIR / "stl_daily_explicit.csv"
CSV_HEATMAP = CSV_DIR / "stl_heatmap_dow_hour.csv"
CSV_DEVICE = DATA_ROOT / "outputs" / "eda_demand_category_device" / "02_device_pct_within_category.csv"
CSV_RANKING = DATA_ROOT / "outputs" / "eda_listing_snapshot" / "_cross" / "01_ranking_by_category.csv"

EVENTS_GLOB = str(DATA_ROOT / "fact_user_events" / "*.parquet")

TET_MUNG1 = pd.Timestamp("2026-02-17")
TET_HOLIDAY_END = pd.Timestamp("2026-02-22")
TET_TREATMENT_DAYS = 12
TET_HOLIDAY_START = TET_HOLIDAY_END - pd.Timedelta(days=TET_TREATMENT_DAYS - 1)

DOW_LABELS = ["T2", "T3", "T4", "T5", "T6", "T7", "CN"]
DOW_ORDER = list(range(7))
HOUR_ORDER = list(range(24))
PEAK_HOURS = list(range(8, 17))
PEAK_COLOR, PEAK_TICK = "#66BB6A", "#43A047"

CAT_ORDER = [1010, 1020, 1030, 1040, 1050]
PALETTE_BAR = {1010: "#52A675", 1020: "#4F8FC9", 1030: "#8E79B8", 1040: "#D66B6B", 1050: "#E09A55"}
DEVICE_COLORS = {"Android": "#1f77b4", "iOS": "#aec7e8", "Desktop": "#d62728", "MSite": "#ff9896"}
DEVICE_EDGE = {"Android": "#1565a8", "iOS": "#8eb8d8", "Desktop": "#b01e22", "MSite": "#e07a7a"}

CSV_DIR.mkdir(parents=True, exist_ok=True)
CHART_DIR.mkdir(parents=True, exist_ok=True)


def _ensure_pptx() -> None:
    try:
        import pptx  # noqa: F401
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])


def _ensure_stl_csv() -> None:
    """One-time DuckDB export → CSV (skip if seeds exist)."""
    if CSV_DAILY.exists() and CSV_HEATMAP.exists():
        return
    try:
        import duckdb
        from statsmodels.tsa.seasonal import STL  # noqa: F401 — only for validation
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "duckdb", "statsmodels", "-q"])
        import duckdb

    login = "is_login = 'login'"
    explicit = "event_type IN ('view_phone', 'contact_chat', 'contact_zalo', 'contact_sms')"
    cats = "1010, 1020, 1030, 1040, 1050"

    con = duckdb.connect(":memory:")
    con.execute("PRAGMA memory_limit='4GB'")
    con.execute("PRAGMA threads=2")

    daily = con.execute(
        f"""
        SELECT CAST(date AS DATE) AS date, COUNT(*)::BIGINT AS y_explicit
        FROM read_parquet('{EVENTS_GLOB}')
        WHERE {login} AND {explicit} AND category IN ({cats})
        GROUP BY 1 ORDER BY 1
        """
    ).df()
    daily.to_csv(CSV_DAILY, index=False)

    heat = con.execute(
        f"""
        SELECT (EXTRACT('dow' FROM event_ts)::INTEGER + 6) % 7 AS dow_mon0,
               EXTRACT('hour' FROM event_ts)::INTEGER AS hour,
               COUNT(*)::BIGINT AS n_explicit
        FROM read_parquet('{EVENTS_GLOB}')
        WHERE {login} AND {explicit} AND category IN ({cats})
        GROUP BY 1, 2
        """
    ).df()
    con.close()
    heat.to_csv(CSV_HEATMAP, index=False)
    print("Exported STL seed CSV:", CSV_DAILY.name, CSV_HEATMAP.name)


def _load_stl_series() -> tuple[pd.Series, pd.Series, pd.Series, dict]:
    from statsmodels.tsa.seasonal import STL

    daily = pd.read_csv(CSV_DAILY, parse_dates=["date"])
    ts = daily.set_index("date")["y_explicit"].asfreq("D").interpolate(method="time")
    mask_tet = (ts.index >= TET_HOLIDAY_START) & (ts.index <= TET_HOLIDAY_END)
    mask_baseline = ~mask_tet

    fit = STL(ts, period=7, robust=True).fit()
    seasonal = pd.Series(fit.seasonal, index=ts.index)
    trend = pd.Series(fit.trend, index=ts.index)

    def dow_profile(s: pd.Series) -> pd.Series:
        return s.groupby(s.index.dayofweek).mean().reindex(range(7))

    prof = dow_profile(seasonal[mask_baseline])
    heat_raw = pd.read_csv(CSV_HEATMAP)
    heat = (
        heat_raw.pivot(index="dow_mon0", columns="hour", values="n_explicit")
        .reindex(index=DOW_ORDER, columns=HOUR_ORDER)
        .fillna(0)
    )
    heat.index = DOW_LABELS
    heat_share = heat / heat.values.sum() * 100.0
    peak_cell = heat_share.stack().idxmax()
    peak_dow, peak_hour = peak_cell[0], int(peak_cell[1])

    daily_peak = []
    for d in DOW_ORDER:
        sub = heat_raw[heat_raw["dow_mon0"] == d]
        if sub.empty:
            continue
        hsub = sub.set_index("hour")["n_explicit"].reindex(HOUR_ORDER, fill_value=0)
        tot = hsub.sum()
        if tot:
            daily_peak.append(hsub.loc[PEAK_HOURS].sum() / tot * 100)

    insights = {
        "peak_dow_label": DOW_LABELS[int(prof.idxmax())],
        "weekend_vs_mid": float(prof.loc[[5, 6]].mean() - prof.loc[1:4].mean()),
        "peak_traffic": f"{peak_dow} {peak_hour:02d}h",
        "peak_traffic_pct": float(heat_share.loc[peak_dow, peak_hour]),
        "peak_hours_avg": float(np.mean(daily_peak)) if daily_peak else 0.0,
        "peak_hours_total": float(heat_share.loc[:, PEAK_HOURS].values.sum()),
        "n_baseline": int(mask_baseline.sum()),
        "t_range": f"{ts.index.min():%Y-%m-%d} → {ts.index.max():%Y-%m-%d}",
        "prof": prof,
        "heat_share": heat_share,
    }
    return ts, trend, prof, insights


def chart_macro(ts: pd.Series, trend: pd.Series, insights: dict) -> Path:
    fig, ax = plt.subplots(figsize=(12, 4.2))
    ax.plot(ts.index, ts.values, color="#bdc3c7", alpha=0.6, label="$y_t$ (observed)")
    ax.plot(trend.index, trend.values, color="#e74c3c", lw=3, label="$T_t$ (STL trend)")
    ax.axvspan(
        TET_HOLIDAY_START,
        TET_HOLIDAY_END + pd.Timedelta(days=1),
        color="#f1c40f",
        alpha=0.22,
        lw=0,
        label=f"Treatment: Tết ({TET_TREATMENT_DAYS} ngày)",
    )
    ax.axvline(TET_MUNG1, color="#d35400", ls="--", lw=1.2, alpha=0.85)
    ax.set_title(f"Macro: explicit contact & STL trend ({insights['t_range']})")
    ax.set_ylabel("$y_t$ — explicit contact / ngày (login)")
    ax.grid(axis="both", linestyle="--", alpha=0.35)
    ax.legend(loc="upper right")
    ax.text(TET_MUNG1, ax.get_ylim()[1] * 0.82, "Mùng 1\n2026-02-17", ha="center", color="#d35400", fontsize=10)
    fig.tight_layout()
    out = CHART_DIR / "01_macro_trend.png"
    fig.savefig(out, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return out


def chart_dow(prof: pd.Series, n_baseline: int) -> Path:
    y_pad = 1.28
    y_max, y_min = prof.max() * y_pad, prof.min() * y_pad
    fig, ax = plt.subplots(figsize=(7, 4.2))
    colors = ["#2E7D32" if v >= 0 else "#C62828" for v in prof.values]
    bars = ax.bar(DOW_LABELS, prof.values, color=colors, alpha=0.85)
    ax.axhline(0, color="black", lw=1.5, linestyle="--")
    pad = (y_max - y_min) * 0.03
    for bar, v in zip(bars, prof.values):
        y_txt, va = (v + pad, "bottom") if v >= 0 else (v - pad, "top")
        ax.text(bar.get_x() + bar.get_width() / 2, y_txt, f"{v:+,.0f}", ha="center", va=va, fontsize=9, fontweight="bold")
    ax.set_title(f"DOW profile — baseline STL $S_t$ ($n={n_baseline}$ ngày)")
    ax.set_ylabel("$S_t$ — seasonal component")
    ax.set_ylim(y_min, y_max)
    ax.grid(axis="y", linestyle="--", alpha=0.35)
    fig.tight_layout()
    out = CHART_DIR / "02_dow_baseline.png"
    fig.savefig(out, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return out


def chart_heatmap(heat_share: pd.DataFrame, insights: dict) -> Path:
    g0, g1 = min(PEAK_HOURS), max(PEAK_HOURS)
    fig, ax = plt.subplots(figsize=(10, 4.5))
    sns.heatmap(
        heat_share,
        ax=ax,
        cmap="YlOrRd",
        linewidths=0.3,
        linecolor="white",
        cbar_kws={"label": "% tổng explicit contact"},
        vmin=0,
    )
    for x in (g0 - 0.5, g1 + 0.5):
        ax.axvline(x=x, color=PEAK_COLOR, linewidth=1.8, linestyle=(0, (6, 4)), alpha=0.85, zorder=10)
    for tick in ax.get_xticklabels():
        try:
            h = int(tick.get_text())
        except ValueError:
            continue
        if g0 <= h <= g1:
            tick.set_color(PEAK_TICK)
            tick.set_fontweight("bold")
    ax.set_title(
        f"Lưu lượng liên hệ — DOW × giờ ({insights['t_range']})\n"
        f"Peak: {insights['peak_traffic']} ({insights['peak_traffic_pct']:.2f}%) · "
        f"Khung {g0}h–{g1}h: {insights['peak_hours_avg']:.1f}%/ngày TB · "
        f"{insights['peak_hours_total']:.1f}% tổng kỳ",
        fontsize=11,
    )
    ax.set_xlabel("Hour (0–23)")
    ax.set_ylabel("Day of week")
    fig.tight_layout()
    out = CHART_DIR / "03_heatmap_dow_hour.png"
    fig.savefig(out, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return out


def chart_device_from_csv() -> Path:
    pivot = pd.read_csv(CSV_DEVICE).set_index("category_name")
    dev_order = [c for c in ["Android", "iOS", "Desktop", "MSite"] if c in pivot.columns]
    pivot = pivot[dev_order]

    def bar_label_color(hex_color: str) -> str:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return "white" if (0.299 * r + 0.587 * g + 0.114 * b) < 140 else "#1a1a1a"

    fig, ax = plt.subplots(figsize=(10.5, 5))
    fig.patch.set_facecolor("#fafafa")
    ax.set_facecolor("#fafafa")
    y = np.arange(len(pivot))
    left = np.zeros(len(pivot))
    for col in pivot.columns:
        vals = pivot[col].to_numpy()
        bars = ax.barh(
            y,
            vals,
            left=left,
            height=0.62,
            label=col,
            color=DEVICE_COLORS.get(col, "#9ecae1"),
            edgecolor=DEVICE_EDGE.get(col, "#7a9eb8"),
            linewidth=0.9,
            alpha=0.92,
        )
        for i, (bar, v) in enumerate(zip(bars, vals)):
            if v >= 3.5:
                ax.text(
                    left[i] + v / 2,
                    bar.get_y() + bar.get_height() / 2,
                    f"{v:.1f}%",
                    ha="center",
                    va="center",
                    fontsize=8.5,
                    fontweight="600",
                    color=bar_label_color(DEVICE_COLORS.get(col, "#9ecae1")),
                )
        left += vals
    ax.set_yticks(y)
    ax.set_yticklabels(pivot.index.tolist())
    ax.invert_yaxis()
    ax.set_xlim(0, 100)
    ax.set_xlabel("Tỉ lệ demand user theo device (%) — mỗi category = 100%")
    ax.set_title(
        "Demand user theo category (README) × thiết bị\nstacked 100% · APPROX_COUNT_DISTINCT(user_id)",
        fontweight="bold",
    )
    ax.grid(axis="x", alpha=0.35, linestyle="--")
    ax.legend(title="Thiết bị", bbox_to_anchor=(1.02, 1), loc="upper left", fontsize=9)
    fig.tight_layout()
    out = CHART_DIR / "04_device_stacked.png"
    fig.savefig(out, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return out


def chart_contact_rate_from_csv() -> Path:
    r = pd.read_csv(CSV_RANKING)
    r["category"] = pd.Categorical(r["category"], categories=CAT_ORDER, ordered=True)
    r = r.sort_values("category")
    r["total_contacts"] = r["avg_contacts_24h"] * r["listing_days"]
    r["share_contacts_pct"] = 100 * r["total_contacts"] / r["total_contacts"].sum()
    n_days = int(r["listing_days"].sum())

    fig, ax = plt.subplots(figsize=(10, 5.5))
    fig.patch.set_facecolor("#FAFAFA")
    ax.set_facecolor("#FAFAFA")
    fig.suptitle("So sánh 5 category — Tỷ lệ liên hệ tổng thể", fontweight="bold", fontsize=13, y=0.97)
    fig.text(0.5, 0.885, f"fact_listing_snapshot · {n_days:,} listing-days", ha="center", fontsize=9, color="#555")
    fig.text(0.5, 0.055, "Công thức: Σcontacts ÷ Σviews (theo category)", ha="center", fontsize=8.5, style="italic", color="#666")

    y = r["category_label"].astype(str)
    colors = [PALETTE_BAR[int(c)] for c in r["category"]]
    bars = ax.barh(y, r["contact_rate_pct"], color=colors, edgecolor="white", height=0.68, alpha=0.92)
    xmax = r["contact_rate_pct"].max() * 1.12
    for bar, v in zip(bars, r["contact_rate_pct"]):
        ax.text(v + xmax * 0.02, bar.get_y() + bar.get_height() / 2, f"{v:.2f}%", va="center", fontsize=9.5, fontweight="600")
    ax.set_xlabel("Tỷ lệ liên hệ (%)")
    ax.set_xlim(0, xmax)
    ax.invert_yaxis()
    ax.grid(axis="x", linestyle="--", alpha=0.35)
    plt.subplots_adjust(top=0.76, bottom=0.20, left=0.28, right=0.96)
    out = CHART_DIR / "05_contact_rate.png"
    fig.savefig(out, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return out


def chart_market_share_pie_from_csv() -> Path:
    r = pd.read_csv(CSV_RANKING)
    r["total_contacts"] = r["avg_contacts_24h"] * r["listing_days"]
    r["share_contacts_pct"] = 100 * r["total_contacts"] / r["total_contacts"].sum()
    r = r.sort_values("share_contacts_pct", ascending=False)

    fig, ax = plt.subplots(figsize=(8, 6))
    colors = [PALETTE_BAR[int(c)] for c in r["category"]]
    wedges, _, autotexts = ax.pie(
        r["share_contacts_pct"],
        colors=colors,
        autopct=lambda pct: f"{pct:.1f}%" if pct >= 3 else "",
        startangle=90,
        counterclock=False,
        pctdistance=0.72,
        wedgeprops={"linewidth": 0.8, "edgecolor": "white"},
    )
    for t in autotexts:
        t.set_fontsize(10)
        t.set_fontweight("bold")
    legend_lines = [
        f"{lbl.split(' — ')[-1]}: {sz:.1f}% · LH {rt:.2f}%"
        for lbl, sz, rt in zip(r["category_label"], r["share_contacts_pct"], r["contact_rate_pct"])
    ]
    ax.legend(wedges, legend_lines, title="Thị phần contacts", loc="upper center", bbox_to_anchor=(0.5, -0.06), fontsize=9, frameon=False)
    fig.suptitle("Thị phần liên hệ theo category (Σ contacts_24h)", fontweight="bold", fontsize=13)
    fig.tight_layout()
    out = CHART_DIR / "06_market_share_pie.png"
    fig.savefig(out, dpi=130, bbox_inches="tight")
    plt.close(fig)
    return out


def build_pptx(charts: dict[str, Path], insights: dict) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    NAVY, GRAY, ACCENT = RGBColor(0x1A, 0x23, 0x7E), RGBColor(0x55, 0x55, 0x55), RGBColor(0xE7, 0x4C, 0x3C)

    def blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def title(slide, t: str, sub: str = "", y=0.35):
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.1), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.text, p.font.size, p.font.bold, p.font.color.rgb = t, Pt(28), True, NAVY
        if sub:
            p2 = tb.text_frame.add_paragraph()
            p2.text, p2.font.size, p2.font.color.rgb = sub, Pt(14), GRAY

    def bullets(slide, items: list[str], top=5.6, size=11):
        tb = slide.shapes.add_textbox(Inches(0.65), Inches(top), Inches(12.0), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size, p.font.color.rgb = Pt(size), GRAY

    def pic(slide, path: Path, left, top, width):
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))

    def callout(slide, text: str, left, top):
        sh = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(4.5), Inches(0.9))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
        sh.line.color.rgb = RGBColor(0xF1, 0xC4, 0x0F)
        p = sh.text_frame.paragraphs[0]
        p.text, p.font.size, p.font.bold, p.font.color.rgb = text, Pt(10), True, ACCENT

    s = blank()
    title(s, "Demand-side EDA — Macro, thời gian & phân khúc", f"{insights['t_range']} · CSV-driven deck")
    bullets(s, ["eda_weekly_stl_tet_regime · eda_demand_category_device · eda_listing_snapshot_cross_pie"], top=4.2, size=12)

    s = blank()
    title(s, "1. Macro — Explicit contact & STL trend", "Shock cấu trúc quanh Tết 2026")
    pic(s, charts["macro"], 0.45, 1.35, 12.4)
    callout(s, "T_t giảm mạnh 11–22/02. Không extrapolate trend pre-Tết vào ngày nghỉ.", 8.0, 5.8)
    bullets(s, ["yt = explicit contact/ngày (login, cat 1010–1050)", "STL period=7, robust=True"])

    s = blank()
    title(s, "2. Mùa vụ tuần & khung giờ", "Baseline loại trừ 12 ngày Tết")
    pic(s, charts["dow"], 0.45, 1.25, 5.9)
    pic(s, charts["heatmap"], 6.55, 1.25, 6.3)
    bullets(
        s,
        [
            f"Peak DOW (St): {insights['peak_dow_label']} · cuối tuần vs mid-week ≈ {insights['weekend_vs_mid']:+.0f}",
            f"Peak traffic: {insights['peak_traffic']} ({insights['peak_traffic_pct']:.2f}%)",
        ],
        top=6.35,
    )

    s = blank()
    title(s, "3. Category × thiết bị", CSV_DEVICE.name)
    rows = [
        ("1010", "Căn Hộ / CC", "MSite 41% · Desktop 28%"),
        ("1020", "Nhà ở", "MSite 52%"),
        ("1030", "VP / MBKD", "Mobile ~42%"),
        ("1040", "Đất", "Android+Desktop ~54%"),
        ("1050", "Phòng trọ", "iOS 24% · MSite 41%"),
    ]
    tbl = s.shapes.add_table(6, 3, Inches(0.5), Inches(1.2), Inches(4.6), Inches(2.6)).table
    for c, h in enumerate(["Mã", "Category", "Device mix"]):
        tbl.cell(0, c).text = h
        tbl.cell(0, c).text_frame.paragraphs[0].font.bold = True
    for r, row in enumerate(rows, 1):
        for c, v in enumerate(row):
            tbl.cell(r, c).text = v
    pic(s, charts["device"], 5.5, 1.15, 7.3)
    bullets(s, ["Mobile đa số; Desktop cao 1010/1040", "MSite lớn 1010/1020/1050 → ưu tiên mobile web"])

    s = blank()
    title(s, "4. Supply snapshot — CVR & thị phần", CSV_RANKING.name)
    pic(s, charts["cvr"], 0.45, 1.2, 6.5)
    pic(s, charts["pie"], 7.1, 1.15, 5.8)
    bullets(s, ["1040 Đất CVR ~11% cao nhất", "1020 ~51% thị phần contacts", "1050 CVR ~7% thấp nhất"])

    s = blank()
    title(s, "Kết luận", "Gợi ý hành động")
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.8), Inches(4.8))
    tf = tb.text_frame
    for i, line in enumerate(
        [
            "Dự báo demand: prior hour×DOW; regime Tết riêng.",
            "Vận hành: CRM 8h–16h, T2–T5; kỳ vọng thấp cuối tuần.",
            "Sản phẩm: MSite 1010/1020/1050; Chat cho user login.",
            "Category: 1020 volume, 1040 CVR — chiến lược lead khác nhau.",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i + 1}. {line}"
        p.font.size, p.font.color.rgb = Pt(16), GRAY

    prs.save(PPTX_PATH)
    return PPTX_PATH


def main() -> None:
    plt.rcParams.update({"figure.facecolor": "white", "axes.titleweight": "bold"})
    _ensure_pptx()
    _ensure_stl_csv()

    ts, trend, prof, insights = _load_stl_series()
    charts = {
        "macro": chart_macro(ts, trend, insights),
        "dow": chart_dow(prof, insights["n_baseline"]),
        "heatmap": chart_heatmap(insights["heat_share"], insights),
        "device": chart_device_from_csv(),
        "cvr": chart_contact_rate_from_csv(),
        "pie": chart_market_share_pie_from_csv(),
    }
    print("Charts:", {k: str(v) for k, v in charts.items()})
    path = build_pptx(charts, insights)
    print("PPTX:", path)


if __name__ == "__main__":
    main()
