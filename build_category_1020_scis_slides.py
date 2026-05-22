"""Generate 2 SCIS slides (HTML) for category 1020 — layout like 1040 Đất."""

from pathlib import Path

DATA_ROOT = Path(__file__).resolve().parent
OUT = DATA_ROOT / "outputs" / "category_1020_business" / "slides"

CSS = """
* { box-sizing: border-box; margin: 0; padding: 0; }
body { font-family: 'Segoe UI', Arial, sans-serif; background: #e8ecf1; }
.slide {
  width: 1280px; height: 720px; margin: 24px auto; background: #fff;
  display: flex; flex-direction: column; overflow: hidden;
  box-shadow: 0 4px 24px rgba(0,0,0,.15);
}
.hdr {
  background: linear-gradient(90deg, #0d47a1 0%, #1565c0 100%);
  color: #fff; padding: 10px 20px 8px;
}
.hdr h1 { font-size: 22px; font-weight: 700; letter-spacing: 0.3px; }
.hdr .sub { font-size: 11px; opacity: 0.92; margin-top: 4px; }
.kpis {
  display: grid; grid-template-columns: repeat(6, 1fr);
  gap: 6px; padding: 8px 12px; background: #f5f7fa;
  border-bottom: 1px solid #dde3eb;
}
.kpi {
  background: #fff; border: 1px solid #cfd8dc;
  border-top: 4px solid var(--accent, #1565c0);
  padding: 6px 8px; text-align: center;
}
.kpi .val { font-size: 20px; font-weight: 800; color: var(--accent, #1565c0); line-height: 1.1; }
.kpi .lbl { font-size: 9px; color: #455a64; margin-top: 4px; line-height: 1.25; }
.body { flex: 1; padding: 8px 14px 6px; overflow: hidden; }
.sec-title {
  font-size: 13px; font-weight: 800; letter-spacing: 0.5px;
  margin: 6px 0 5px; padding-bottom: 3px;
  border-bottom: 2px solid var(--sec-color, #1565c0);
  color: var(--sec-color, #1565c0);
}
.sec-title.red { --sec-color: #c62828; color: #c62828; border-color: #c62828; }
.row { display: flex; gap: 8px; margin-bottom: 6px; min-height: 0; }
.label-box {
  width: 168px; flex-shrink: 0;
  background: #263238; color: #fff;
  border-radius: 6px; padding: 10px 10px;
  font-size: 11px; font-weight: 700; line-height: 1.35;
  display: flex; align-items: center;
}
.content-box {
  flex: 1; border: 2px solid #90caf9; border-radius: 4px;
  padding: 8px 12px; font-size: 11px; line-height: 1.45; color: #212121;
}
.content-box b { color: #0d47a1; }
.content-box .hi { color: #c62828; font-weight: 700; }
.content-box ol { margin: 4px 0 0 18px; }
.content-box li { margin-bottom: 3px; }
.footer {
  background: #0d47a1; color: #fff;
  display: flex; justify-content: space-between; align-items: center;
  padding: 6px 20px; font-size: 11px; font-weight: 600;
}
.footer .active { text-decoration: underline; }
.footer .page { opacity: 0.9; }
"""

SLIDE1_BODY = """
<div class="kpis">
  <div class="kpi" style="--accent:#1565c0"><div class="val">1,51M</div><div class="lbl">listings catalog<br>(~2,5× cat 1010)</div></div>
  <div class="kpi" style="--accent:#c62828"><div class="val">72%</div><div class="lbl">sell — DNA<br>thị trường bán nhà</div></div>
  <div class="kpi" style="--accent:#2e7d32"><div class="val">20,5% / 19,3%</div><div class="lbl">CVR let / sell<br>(let &gt; sell, ngược 1010)</div></div>
  <div class="kpi" style="--accent:#ef6c00"><div class="val">19% / 8,6%</div><div class="lbl">contact-day let / sell<br>(bán = nút thắt)</div></div>
  <div class="kpi" style="--accent:#6a1b9a"><div class="val">6.187</div><div class="lbl">tin HQ underexposed<br>(~2,7× vs 1010)</div></div>
  <div class="kpi" style="--accent:#00838f"><div class="val">~70%</div><div class="lbl">supply TP.HCM<br>micro-market</div></div>
</div>
<div class="body">
  <div class="sec-title">SITUATION</div>
  <div class="row">
    <div class="label-box">Thị trường<br>Bán nhà scale</div>
    <div class="content-box">
      <b>1,51M tin</b>, <b>72% sell</b> — đây là transaction-heavy house market, không phải rental-first như căn hộ.
      Cấu trúc tin: <b>house_type 100%</b>, floors 61%, width 76% → matching phức tạp hơn 1010.
      <b>~270k login users</b> (sample 6–8%) — volume demand lớn nhất trong cặp 1010/1020.
    </div>
  </div>
  <div class="row">
    <div class="label-box">CVR &amp; Snapshot<br>đảo chiều 1010</div>
    <div class="content-box">
      Catalog CVR: <b>let 20,45%</b> &gt; <b>sell 19,32%</b> (+1,1pp) — <b>ngược 1010</b> (sell &gt; let).
      Snapshot thực tế: <b>cho thuê nhà 19% ngày có contact</b> vs <b>bán chỉ 8,55%</b> trên ~7,5M listing-days.
      → KPI và message phải <b>tách let vs sell</b>, không gộp “Nhà ở” chung một funnel.
    </div>
  </div>
  <div class="sec-title red">CHALLENGES</div>
  <div class="row">
    <div class="label-box">3 Rào cản<br>Cốt lõi</div>
    <div class="content-box">
      <ol>
        <li><b>KPI trap:</b> chỉ <b>8,38% session</b> có explicit contact; <b>50,8% events</b> là other_interaction (<b>96% ad_view</b>) — đo nhầm exposure = lead.</li>
        <li><b>Gem bị chôn:</b> <b>~6,2k</b> HQ underexposed (271% contact/PV, exposure=2) vs <b>~4,3k</b> oversaturated (58%, exposure median <b>13</b>).</li>
        <li><b>Weak pocket sell:</b> <b>30–50m² ~310k tin</b>, CVR <b>17,5%</b> — bucket lớn nhất, kéo conversion category bán.</li>
      </ol>
    </div>
  </div>
  <div class="row">
    <div class="label-box">2 Phát hiện<br>mới</div>
    <div class="content-box">
      <b>Broker/spam user 12,38%</b> — cao hơn 1010 (9,6%); high-intent chỉ <b>5,98%</b>.<br>
      <b>Bulk seller cluster:</b> ~<b>102 listings/seller</b>, chỉ <b>0,29 PV/listing</b> — cold/spam supply làm loãng feed.<br>
      <b>Let mặt phố ~15% CVR</b> — pocket thuê yếu riêng (khác sell ngõ/mặt phố).
    </div>
  </div>
</div>
"""

SLIDE2_BODY = """
<div class="kpis">
  <div class="kpi" style="--accent:#2e7d32"><div class="val">+93k</div><div class="lbl">contact events<br>(boost underexposed)</div></div>
  <div class="kpi" style="--accent:#1565c0"><div class="val">2,7×</div><div class="lbl">opportunity vs<br>cat 1010 boost</div></div>
  <div class="kpi" style="--accent:#ef6c00"><div class="val">+9,7k</div><div class="lbl">sessions w/ contact<br>(explicit +1pp)</div></div>
  <div class="kpi" style="--accent:#6a1b9a"><div class="val">~322k</div><div class="lbl">imp-equiv saved<br>(demote 4,3k tin)</div></div>
  <div class="kpi" style="--accent:#c62828"><div class="val">~6,2k</div><div class="lbl">listings thêm contact<br>(sell 30–50m² +2pp)</div></div>
  <div class="kpi" style="--accent:#00838f"><div class="val">4,9k</div><div class="lbl">sell HCM trong<br>segment boost</div></div>
</div>
<div class="body">
  <div class="sec-title">STRATEGIES</div>
  <div class="row">
    <div class="label-box">Quick Wins<br>Tháng 1–2</div>
    <div class="content-box">
      ① <b>Tách KPI layer</b> — North Star = explicit_contact/session; loại ad_view khỏi báo cáo conversion.<br>
      ② <b>Boost 6,2k HQ underexposed</b> — feed rank từ <code>10_health_ranked_underexposed_1020.csv</code>; ưu tiên <b>sell HCMC (~4,9k tin)</b>.<br>
      ③ <b>Phone-first CTA</b> — 77% explicit = view_phone; nút gọi/Zalo nổi trên listing sell.
    </div>
  </div>
  <div class="row">
    <div class="label-box">Mid-term<br>Tháng 2–4</div>
    <div class="content-box">
      ① <b>Demote 4,3k oversaturated</b> — cap impression (exposure median 13, cao nhất trong cặp cat).<br>
      ② <b>Completeness gate</b> — bắt buộc house_type + floors + legal + width trước rank; ưu tiên ngõ/mặt phố sell.<br>
      ③ <b>Deep-compare UX</b> — 60% session có deep_compare; checklist so sánh trước khi contact (journey mua nhà dài).
    </div>
  </div>
  <div class="row">
    <div class="label-box">Full<br>6 tháng</div>
    <div class="content-box">
      ① <b>Android mobile-first</b> (26% events) — form đăng tin &amp; xem listing tối ưu mobile.<br>
      ② <b>Cap bulk seller</b> — giới hạn listings/seller; onboarding seller chất lượng thay spam inventory.<br>
      ③ <b>Fix pocket sell 30–50m²</b> — pricing/ảnh/định vị review; mục tiêu CVR 17,5% → 19,5% trên ~310k tin.
    </div>
  </div>
  <div class="sec-title red">IMPACTS</div>
  <div class="row">
    <div class="label-box">Recommender<br>Hooks</div>
    <div class="content-box">
      <b>Boost underexposed:</b> exposure 2→4 trên 6,187 tin → <b>+~93k contact events</b> (ước tính sample).<br>
      <b>Demote oversaturated:</b> giải phóng <b>~322k impression-equiv</b> — chuyển inventory sang gem &amp; normal segment.<br>
      <b>Session +1pp explicit</b> (8,38→9,38%): <b>+~9,7k sessions</b> có contact trong cohort login.
    </div>
  </div>
  <div class="row">
    <div class="label-box">Root Cause<br>&amp; Kết luận</div>
    <div class="content-box">
      <b>Vấn đề gốc:</b> Feed ưu tiên <span class="hi">visibility (ad_view)</span> hơn <span class="hi">verified converters</span>; sell volume lớn nhưng contact-day thấp.<br>
      <b>Kết luận:</b> 1020 = <b>đại dương bán nhà</b> — lever #1 là recsys boost gem (scale gấp 3× căn hộ); lever #2 là fix sell 30–50m² + dọn bulk seller.<br>
      <b>Caveat:</b> explicit ≠ ad_view · CVR catalog ≠ clustering cohort · sizing từ sample 6–8%.
    </div>
  </div>
</div>
"""


def _html(title: str, subtitle: str, body: str, footer_left: str, footer_right: str, page: str) -> str:
    return f"""<!DOCTYPE html>
<html lang="vi"><head><meta charset="utf-8"><title>{title}</title>
<style>{CSS}</style></head><body>
<div class="slide">
  <div class="hdr"><h1>{title}</h1><div class="sub">{subtitle}</div></div>
  {body}
  <div class="footer"><span class="{footer_left[1]}">{footer_left[0]}</span>
    <span class="{footer_right[1]}">{footer_right[0]}</span><span class="page">{page}</span></div>
</div></body></html>"""


def build_pptx() -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    NAVY = RGBColor(0x0D, 0x47, 0xA1)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x26, 0x32, 0x38)
    BLUE = RGBColor(0x15, 0x65, 0xC0)
    RED = RGBColor(0xC6, 0x28, 0x28)
    BODY = RGBColor(0x21, 0x21, 0x21)
    GRAY = RGBColor(0x45, 0x5A, 0x64)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def fill_rect(slide, l, t, w, h, rgb):
        sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb
        sh.line.fill.background()
        return sh

    def textbox(slide, l, t, w, h, lines: list[tuple[str, int, bool, object]], wrap=True):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        for i, (txt, size, bold, color) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = txt
            p.font.size, p.font.bold, p.font.color.rgb = Pt(size), bold, color
        return tb

    def header(slide, title: str, sub: str):
        fill_rect(slide, 0, 0, 13.333, 0.72, NAVY)
        textbox(slide, 0.35, 0.12, 12.6, 0.55, [(title, 20, True, WHITE)])
        textbox(slide, 0.35, 0.48, 12.6, 0.22, [(sub, 9, False, WHITE)])

    def footer(slide, active: str):
        fill_rect(slide, 0, 7.05, 13.333, 0.45, NAVY)
        textbox(
            slide,
            0.4,
            7.12,
            12.5,
            0.3,
            [
                ("SITUATION & CHALLENGES", 10, active == "s1", WHITE),
                ("STRATEGIES & IMPACTS", 10, active == "s2", WHITE),
            ],
        )
        tb = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.text = "1 / 2" if active == "s1" else "2 / 2"
        p.font.size, p.font.color.rgb, p.alignment = Pt(10), WHITE, PP_ALIGN.RIGHT

    def kpis(slide, items: list[tuple[str, str, object]], top=0.78):
        w = 2.1
        for i, (val, lbl, accent) in enumerate(items):
            l = 0.25 + i * (w + 0.05)
            fill_rect(slide, l, top, w, 0.05, accent)
            fill_rect(slide, l, top + 0.05, w, 0.62, WHITE)
            sh = slide.shapes.add_shape(1, Inches(l), Inches(top + 0.05), Inches(w), Inches(0.62))
            sh.fill.solid()
            sh.fill.fore_color.rgb = WHITE
            sh.line.color.rgb = RGBColor(0xCF, 0xD8, 0xDC)
            textbox(slide, l + 0.05, top + 0.1, w - 0.1, 0.28, [(val, 16, True, accent)])
            textbox(slide, l + 0.05, top + 0.38, w - 0.1, 0.28, [(lbl, 7, False, GRAY)])

    def sec(slide, label: str, y: float, red=False):
        c = RED if red else BLUE
        tb = slide.shapes.add_textbox(Inches(0.35), Inches(y), Inches(3), Inches(0.28))
        p = tb.text_frame.paragraphs[0]
        p.text, p.font.size, p.font.bold, p.font.color.rgb = label, Pt(11), True, c

    def row(slide, label: str, body: str, y: float, h: float):
        fill_rect(slide, 0.35, y, 1.55, h, DARK)
        textbox(slide, 0.42, y + 0.06, 1.4, h - 0.1, [(label, 9, True, WHITE)])
        sh = slide.shapes.add_shape(1, Inches(2.0), Inches(y), Inches(11.0), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = WHITE
        sh.line.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
        textbox(slide, 2.08, y + 0.06, 10.85, h - 0.1, [(body, 8.5, False, BODY)])

    sub = "Chợ Tốt BĐS | Datathon 2026 | Category 1020 — Nhà ở | EDA: performance · behavior · clustering · bridge"
    kpi1 = [
        ("1,51M", "listings catalog\n(~2,5× cat 1010)", BLUE),
        ("72%", "sell — DNA thị trường bán nhà", RED),
        ("20,5% / 19,3%", "CVR let / sell (let > sell)", RGBColor(0x2E, 0x7D, 0x32)),
        ("19% / 8,6%", "contact-day let / sell", RGBColor(0xEF, 0x6C, 0x00)),
        ("6.187", "tin HQ underexposed", RGBColor(0x6A, 0x1B, 0x9A)),
        ("~70%", "supply TP.HCM", RGBColor(0x00, 0x83, 0x8F)),
    ]

    s1 = blank()
    header(s1, "SITUATION & CHALLENGES — CATEGORY 1020 (NHÀ Ở)", sub)
    kpis(s1, kpi1)
    sec(s1, "SITUATION", 1.52)
    row(
        s1,
        "Thị trường\nBán nhà scale",
        "1,51M tin, 72% sell — transaction-heavy house market. house_type 100%, floors 61%, width 76%. ~270k login users (sample 6–8%).",
        1.82,
        0.72,
    )
    row(
        s1,
        "CVR & Snapshot\nđảo chiều 1010",
        "Catalog: let 20,45% > sell 19,32% (+1,1pp) — ngược 1010. Snapshot: cho thuê 19% contact-day vs bán 8,55% trên ~7,5M listing-days. KPI phải tách let vs sell.",
        2.58,
        0.72,
    )
    sec(s1, "CHALLENGES", 3.38, red=True)
    row(
        s1,
        "3 Rào cản\nCốt lõi",
        "1) KPI trap: 8,38% session explicit; 50,8% events other_interaction (96% ad_view). "
        "2) Gem bị chôn: ~6,2k HQ underexposed vs ~4,3k oversaturated (exp. median 13). "
        "3) Weak pocket sell: 30–50m² ~310k tin, CVR 17,5%.",
        3.68,
        0.95,
    )
    row(
        s1,
        "2 Phát hiện\nmới",
        "Broker/spam user 12,38% (cao hơn 1010). Bulk seller ~102 listings/seller, 0,29 PV/listing. Let mặt phố ~15% CVR — pocket thuê yếu.",
        4.68,
        0.72,
    )
    footer(s1, "s1")

    kpi2 = [
        ("+93k", "contact (boost gem)", RGBColor(0x2E, 0x7D, 0x32)),
        ("2,7×", "opportunity vs 1010", BLUE),
        ("+9,7k", "sessions (+1pp explicit)", RGBColor(0xEF, 0x6C, 0x00)),
        ("~322k", "imp-equiv (demote)", RGBColor(0x6A, 0x1B, 0x9A)),
        ("~6,2k", "listings (+2pp sell pocket)", RED),
        ("4,9k", "sell HCM in boost seg.", RGBColor(0x00, 0x83, 0x8F)),
    ]
    sub2 = "Mục tiêu: boost gem + fix sell 30–50m² | Focus: sell HCMC · phone-first · deep-compare"
    s2 = blank()
    header(s2, "STRATEGIES & IMPACTS — CATEGORY 1020 (NHÀ Ở)", sub2)
    kpis(s2, kpi2)
    sec(s2, "STRATEGIES", 1.52)
    row(
        s2,
        "Quick Wins\nTháng 1–2",
        "① Tách KPI layer (explicit_contact/session). ② Boost 6,2k HQ underexposed — ưu tiên sell HCMC (~4,9k). ③ Phone-first CTA (77% explicit = view_phone).",
        1.78,
        0.58,
    )
    row(
        s2,
        "Mid-term\nTháng 2–4",
        "① Demote 4,3k oversaturated (cap impression). ② Completeness gate (house_type, floors, legal, width). ③ Deep-compare UX (60% session).",
        2.4,
        0.58,
    )
    row(
        s2,
        "Full\n6 tháng",
        "① Android mobile-first (26% events). ② Cap bulk seller. ③ Fix pocket sell 30–50m² CVR 17,5% → 19,5%.",
        3.02,
        0.58,
    )
    sec(s2, "IMPACTS", 3.68, red=True)
    row(
        s2,
        "Recommender\nHooks",
        "Boost: exposure 2→4 trên 6,187 tin → +~93k contact events. Demote: ~322k impression-equiv. Session +1pp: +~9,7k sessions có contact.",
        3.98,
        0.62,
    )
    row(
        s2,
        "Root Cause\n& Kết luận",
        "Gốc: feed ưu tiên visibility (ad_view) hơn verified converters; sell volume lớn nhưng contact-day thấp. "
        "Kết luận: lever #1 recsys boost gem (3× căn hộ); lever #2 fix sell 30–50m² + bulk seller. Caveat: explicit ≠ ad_view.",
        4.65,
        0.62,
    )
    footer(s2, "s2")

    out = OUT / "category_1020_SCIS_slides.pptx"
    prs.save(out)
    return out


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    sub = "Chợ Tốt BĐS | Datathon 2026 | Category 1020 — Nhà ở | EDA: performance · behavior · clustering · bridge"

    p1 = OUT / "slide_01_situation_challenges.html"
    p1.write_text(
        _html(
            "SITUATION &amp; CHALLENGES — CATEGORY 1020 (NHÀ Ở)",
            sub,
            SLIDE1_BODY,
            ("SITUATION & CHALLENGES", "active"),
            ("STRATEGIES & IMPACTS", ""),
            "1 / 2",
        ),
        encoding="utf-8",
    )

    p2 = OUT / "slide_02_strategies_impacts.html"
    p2.write_text(
        _html(
            "STRATEGIES &amp; IMPACTS — CATEGORY 1020 (NHÀ Ở)",
            "Mục tiêu: boost gem underexposed + fix sell 30–50m² | Focus: sell HCMC · phone-first · deep-compare",
            SLIDE2_BODY,
            ("SITUATION & CHALLENGES", ""),
            ("STRATEGIES & IMPACTS", "active"),
            "2 / 2",
        ),
        encoding="utf-8",
    )

    index = OUT / "index.html"
    index.write_text(
        f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<title>1020 SCIS Slides</title></head><body style="font-family:sans-serif;padding:24px">
<h1>Category 1020 — SCIS Slides</h1>
<ul>
<li><a href="slide_01_situation_challenges.html">Slide 1 — SITUATION & CHALLENGES</a></li>
<li><a href="slide_02_strategies_impacts.html">Slide 2 — STRATEGIES & IMPACTS</a></li>
</ul>
<p>Mở file HTML trong browser → Print/PDF hoặc chụp màn hình 1280×720.</p>
</body></html>""",
        encoding="utf-8",
    )
    pptx = build_pptx()
    print(f"Wrote:\n  {p1}\n  {p2}\n  {index}\n  {pptx}")


if __name__ == "__main__":
    main()
