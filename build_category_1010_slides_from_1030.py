"""Build 2-slide PPTX for category 1010 by cloning layout from 1030_slides.pptx."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation

DATA_ROOT = Path(__file__).resolve().parent
TEMPLATE = DATA_ROOT / "1030_slides.pptx"
OUT_DIR = DATA_ROOT / "outputs" / "category_1010_business" / "slides"
OUT_PPTX = OUT_DIR / "category_1010_SCIS_slides.pptx"


def _delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def _replace_in_slide(slide, replacements: list[tuple[str, str]]) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        new_text = text
        for old, new in replacements:
            if old in new_text:
                new_text = new_text.replace(old, new)
        if new_text != text:
            shape.text_frame.text = new_text


def _set_shape_text(slide, index: int, text: str) -> None:
    sh = slide.shapes[index]
    if sh.has_text_frame:
        sh.text_frame.text = text


def _apply_slide1(sc) -> None:
    """SITUATION & CHALLENGES — layout from 1030 slide 2."""
    reps = [
        (
            "SITUATION & CHALLENGES — CATEGORY 1030 (VĂN PHÒNG / MẶT BẰNG)",
            "SITUATION & CHALLENGES — CATEGORY 1010 (CĂN HỘ / CHUNG CƯ)",
        ),
        ("93,7%", "612k"),
        ("Supply Let vs\n6,3% Sell", "74% let /\n26% sell"),
        ("92,2%", "50,6%"),
        ("other_interaction\n(dominant B2B)", "event layer\n(94% ad_view)"),
        ("+6,3pp", "7,78%"),
        ("Thiếu hụt cung\n30–50 m²", "session explicit\ncontact only"),
        ("76,1%", "74,1%"),
        ("Supply là Agent", "supply TP.HCM\n(project×PN)"),
        ("4,7%", "41%"),
        ("legal_status\nNULL", "có project_id"),
        ("1,6×", "11,1%"),
        ("Non-login vs Login\nPR (bẫy đo)", "contact-day let\n(snapshot)"),
        ("Thị trường\nB2B dual", "Thị trường\nThuê căn HCM"),
        (
            "Category duy nhất có Let lẫn Sell đáng kể: startup/SME (30–100 m²) vs doanh nghiệp (>200 m²).\n"
            "Positive 52,5% · Revenue CVR 4,1% — nhu cầu B2B ổn định, khác residential browsing.",
            "611.823 listings · 74,3% let / 25,7% sell — cho thuê căn hộ tại TP.HCM.\n"
            "41% project_id · bedrooms · ~154k login users (sample 6–8%).",
        ),
        ("Địa lý\n& BCG", "Sell vs Let\nCVR"),
        (
            "Supply tập trung TP.HCM + Hà Nội; tier-2 (Đà Nẵng, Bình Dương…) ít tin nhưng CVR cao → Stars.\n"
            "C5: tỉnh tier-2 outperform trên mỗi listing — cơ hội seller acquisition ngoài HCM+HN.",
            "Catalog CVR: sell 23,0% > let 18,8% (+4,2pp) — bán căn contact tốt hơn thuê.\n"
            "Snapshot thuê: 11,07% contact-day · tuổi tin ~52 ngày.\n"
            "→ KPI & CRM tách theo ad_type + thành phố.",
        ),
        ("Kênh &\nSeller", "Funnel &\nKênh lead"),
        (
            "Agent 76,1% supply · PR 54,0% vs Private 51,0% — khác nhà ở; cần playbook riêng cho B2B.\n"
            "Kênh: other_interaction 92,2% · view_phone 6,4% · chat 1,1% — hành vi so sánh trước khi liên hệ.",
            "7,78% session explicit contact · 50,6% events other_interaction (94% ad_view).\n"
            "High-intent 7,68% · broker/spam 9,58% · 30% explicit = contact_chat.",
        ),
        ("C1 AREA", "C1 KPI"),
        (
            "30–50 m² thiếu cung.\n<30 m² +6,2pp.\n>200 m² thừa −12,6pp.",
            "7,78% session\nexplicit only.\n50,6% ad_view.",
        ),
        ("C2 SCHEMA", "C2 GEM"),
        ("NULL cao", "2.328"),
        (
            "ward/direction/\nproject_id thưa.\nlegal ~4,7% NULL.",
            "304% contact/PV\nexposure = 2.\nBoost priority.",
        ),
        ("C3 AGE", "C3 POCKET"),
        ("↓14 ngày", "289k"),
        (
            "Contacts/views\ngiảm sau 2 tuần.\nCần freshness boost.",
            "1-bed let ~18% CVR.\nVolume lớn nhất.",
        ),
        ("C4 PRICE", "C4 DEMOTE"),
        ("Gap", "2.092"),
        (
            "Bucket giá\nlệch cung/cầu\n(theo segment).",
            "Oversaturated\nexp. median 8.\nCap impression.",
        ),
        ("C5 GEO", "C5 DATA"),
        ("HCM+HN", "~85%"),
        (
            "Tier-2 CVR cao\nsupply thấp.\nMở rộng địa lý.",
            "Tin pre-EDA window\nlàm méo CVR catalog.\nFilter in-window.",
        ),
        ("C6 LET/SELL", "C6 UX"),
        ("94/6", "30%"),
        (
            "Hai thị trường\nmột feed — Sell\nbị chôn vùi.",
            "Chat-first CTA +\nDesktop/iOS 73%.\nNight push ~500 user.",
        ),
        (
            "Root cause: thiếu đúng phân khúc (30–100 m²) · schema thương mại yếu · Let/Sell chung feed · không tách B2B lifecycle.",
            "Root cause: feed ưu tiên ad_view hơn explicit lead · gem underexposed bị chôn · 1PN let kéo CVR · KPI không tách let/sell.",
        ),
        ("IMPACTS & ROADMAP", "STRATEGIES & IMPACTS"),
        ("2 / 4", "1 / 2"),
    ]
    _replace_in_slide(sc, reps)
    _set_shape_text(sc, 13, "18,8%")
    _set_shape_text(sc, 14, "23,0%\nCVR catalog\nsell > let")
    _set_shape_text(sc, 45, "7,78%")
    _set_shape_text(
        sc,
        40,
        "7,78% session explicit · 50,6% other_interaction (94% ad_view).\n"
        "30% explicit = contact_chat · Desktop 40% + iOS 33% — UX thuê căn hộ.",
    )
    _set_shape_text(
        sc,
        52,
        "304% contact/pageview\nexposure = 2.\nBoost recsys rank.",
    )


def _apply_slide2(sc) -> None:
    """STRATEGIES & IMPACTS — layout from 1030 slide 4."""
    reps = [
        (
            "IMPACTS & ROADMAP — CATEGORY 1030 (VĂN PHÒNG / MẶT BẰNG)",
            "STRATEGIES & IMPACTS — CATEGORY 1010 (CĂN HỘ / CHUNG CƯ)",
        ),
        (
            "Mục tiêu: +25% CVR @ W4 · +45% @ M3 · +80% @ M12",
            "Mục tiêu: boost 2,3k gem · demote 2,1k · fix 1PN let | Focus: let HCM + sell project",
        ),
        ("+25%", "+39k"),
        ("CVR lift\nTuần 1–4", "contact events\n(boost gem)"),
        ("+45%", "+3,7k"),
        ("CVR lift\nTháng 2–3", "sessions w/\ncontact (+1pp)"),
        ("+80%", "+68k"),
        ("CVR lift\nTháng 4–12", "imp-equiv saved\n(demote)"),
        ("+42%", "+5,8k"),
        ("Tier-2 geo\nuplift model", "listings +\ncontact (1PN)"),
        ("+15%", "2,3k"),
        ("Freshness\n0–14 ngày", "underexposed\nin boost seg."),
        ("Fill", "30%"),
        ("legal_status\n~4,7% NULL", "explicit via\ncontact_chat"),
        ("ROADMAP", "STRATEGIES"),
        ("Quick Wins (Tuần 1–4) → +25%", "Quick Wins (Tuần 1–2)"),
        (
            "① Bổ sung supply 50–100 m² (Let)\n② Hoàn thiện Legal Status (form + data fill)\n③ Image Quality Gate (≥5 ảnh)",
            "① Tách KPI — explicit_contact/session (không gộp ad_view)\n② Boost 2.328 HQ underexposed (recsys rank)\n③ Chat-first CTA trên listing thuê",
        ),
        ("Mid-term (Tháng 2–3) → +45%", "Mid-term (Tháng 2–4)"),
        (
            "④ Freshness Ranking (ưu tiên 0–14 ngày)\n⑤ Let vs Sell channel separation\n⑥ Tier-2 city listing expansion",
            "④ Demote 2.092 oversaturated — cap impression share\n⑤ Desktop + iOS UX (40% + 33% events)\n⑥ Night-active push (~500 users, 99% night)",
        ),
        ("Full Build (Tháng 4–12) → +80%", "Full (6 tháng)"),
        (
            "⑦ Legal Verification Trust Badge\n⑧ Sell Market Activation (HN + Đà Nẵng)\n⑨ Commercial Relevance Scoring v2 (ML)",
            "⑦ 1PN let CVR 18→20% (completeness + furnishing)\n⑧ Sell: project 2PN, 3–4+ PN quality gate\n⑨ health_segment_boost v2 (let|sell × city)",
        ),
        (
            "1. Mid-size supply (30–100 m²): lấp gap area lớn nhất — CAO · Q1\n2. Tier-2 expansion: +42% CVR model → +XK revenue events — TB-CAO · Q2\n3. Freshness + legal fill + Sell activation — cumulative +80% @ M12",
            "1. Boost underexposed: +~39,5k contact events (exposure 2→4)\n2. Demote oversaturated: ~68k impression-equiv freed\n3. Session +1pp explicit: +~3,7k sessions · 1PN +2pp: ~5,8k listings",
        ),
        (
            "commercial_relevance_1030 · ad_type_channel (let|sell) · area_price_segment · legal_trust_score · freshness_boost\nKhông copy playbook 1050/1040 — B2B cần tách kênh & segment theo diện tích×giá.",
            "health_segment_boost · ad_type_channel (let|sell) · hcm_city_priority · chat_first_score · project_bedroom_gate · night_push_flag\nLever #1 recsys gem 2,3k; lever #2 pocket 1PN let 289k tin.",
        ),
        (
            "Tuần 1: campaign seller 30–50 & 50–100 m² · bắt buộc legal_status + ≥5 ảnh · prototype tách Let/Sell feed.\nFocus tier-2: Đà Nẵng, Bình Dương — CVR cao, supply thấp (BCG Stars).",
            "Tuần 1–2: boost 2,3k gem + demote pilot · chat-first trên listing let.\nFocus HCM (~74% supply) · sell ưu tiên project-linked 2–4 PN. Caveat: explicit ≠ ad_view · sample 6–8%.",
        ),
        ("4 / 4", "2 / 2"),
    ]
    _replace_in_slide(sc, reps)
    _set_shape_text(
        sc,
        46,
        "1. Boost 2,3k underexposed: +~39,5k contact events — CAO · Tuần 1–2\n"
        "2. Demote 2,1k oversaturated: ~68k impression-equiv — CAO · Tuần 2–4\n"
        "3. Session +1pp + 1PN let +2pp: +~9,5k incremental contacts — TB-CAO · Q2",
    )
    _set_shape_text(sc, 30, "Quick Wins (Tuần 1–2) → +39k")
    _set_shape_text(sc, 34, "Mid-term (Tháng 2–4) → +3,7k")
    _set_shape_text(sc, 38, "Full (6 tháng) → +5,8k listings")
    _set_shape_text(sc, 61, "STRATEGIES & IMPACTS")
    _set_shape_text(sc, 63, "STRATEGIES & IMPACTS")


def build() -> Path:
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Missing template: {TEMPLATE}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    shutil.copy2(TEMPLATE, OUT_PPTX)
    prs = Presentation(str(OUT_PPTX))

    # Remove overview (0) and strategies-only (2); keep S&C (1) and Impacts (3)
    _delete_slide(prs, 2)
    _delete_slide(prs, 0)

    _apply_slide1(prs.slides[0])
    _apply_slide2(prs.slides[1])

    prs.save(str(OUT_PPTX))
    return OUT_PPTX


def main() -> None:
    out = build()
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
